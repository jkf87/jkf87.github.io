#!/usr/bin/env python3
"""Append Memory Wiki slides to existing McKinsey PPTX"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Load existing
prs = Presentation('/Users/conanssam-m4/Downloads/openclaw-memory-mckinsey.pptx')
W = prs.slide_width
H = prs.slide_height
blank_layout = prs.slide_layouts[6]

# Palette (same as existing)
NAVY = RGBColor(0x00, 0x2B, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF5, 0xF5, 0xF5)
LIGHT_GREY = RGBColor(0xE8, 0xE8, 0xE8)
MED_GREY = RGBColor(0x99, 0x99, 0x99)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
SIGNAL_RED = RGBColor(0xC8, 0x10, 0x2E)
GOLD = RGBColor(0xB8, 0x86, 0x0B)
TEAL = RGBColor(0x00, 0x7B, 0x7F)
SOFT_BLUE = RGBColor(0x3A, 0x7C, 0xBD)
AMBER = RGBColor(0xD4, 0x8B, 0x0B)
GREEN_ACC = RGBColor(0x1B, 0x7F, 0x4B)
PURPLE = RGBColor(0x6B, 0x21, 0xA8)

TOTAL = 20  # will be 20 total

def set_slide_bg(slide, color=WHITE):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill, line=None, lw=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = lw or Pt(0.5)
    else: s.line.fill.background()
    return s

def tb(slide, l, t, w, h, text, sz=14, c=DARK_GREY, b=False, a=PP_ALIGN.LEFT, fn='Calibri', i=False):
    tx = slide.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = c; p.font.bold = b; p.font.name = fn; p.font.italic = i; p.alignment = a
    return tx

def rich(slide, l, t, w, h, lines, fn='Calibri'):
    """lines: [(text, size, color, bold, align, italic)]"""
    tx = slide.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame; tf.word_wrap = True
    for i, ld in enumerate(lines):
        text = ld[0]; sz = ld[1] if len(ld)>1 else 10; cl = ld[2] if len(ld)>2 else DARK_GREY
        bl = ld[3] if len(ld)>3 else False; al = ld[4] if len(ld)>4 else PP_ALIGN.LEFT
        it = ld[5] if len(ld)>5 else False
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = text; p.font.size = Pt(sz); p.font.color.rgb = cl
        p.font.bold = bl; p.font.name = fn; p.font.italic = it; p.alignment = al; p.space_after = Pt(3)
    return tx

def hdr(slide, title):
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.9), NAVY)
    tb(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.6), title, sz=22, c=WHITE, b=True, fn='Georgia')

def ftr(slide, n):
    add_rect(slide, Inches(0), H-Inches(0.3), W, Inches(0.3), LIGHT_GREY)
    tb(slide, Inches(0.6), H-Inches(0.28), Inches(5), Inches(0.25), "OpenClaw Memory System", sz=8, c=MED_GREY)
    tb(slide, W-Inches(1.5), H-Inches(0.28), Inches(1), Inches(0.25), f"{n}", sz=8, c=MED_GREY, a=PP_ALIGN.RIGHT)

def rule(slide, l, t, w, c=SIGNAL_RED, th=Pt(2)):
    add_rect(slide, l, t, w, th, c)


# ═══════════════════════════════════════════════════════════
# SLIDE 16 — SECTION DIVIDER: MEMORY WIKI
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, NAVY)
rule(slide, Inches(1.2), Inches(2.5), Inches(2), GOLD, Pt(3))
tb(slide, Inches(1.2), Inches(2.8), Inches(10), Inches(1.2),
   "Memory Wiki\n지식 베이스 레이어", sz=40, c=WHITE, b=True, fn='Georgia')
tb(slide, Inches(1.2), Inches(4.3), Inches(10), Inches(0.5),
   "Karpathy의 Belief Layer 철학 × OpenClaw 컴파일 파이프라인 × Obsidian 연동",
   sz=16, c=RGBColor(0xAA, 0xCC, 0xE8))
rule(slide, Inches(1.2), Inches(5.0), Inches(1.5), GOLD, Pt(2))
tb(slide, Inches(1.2), Inches(5.3), Inches(10), Inches(0.4),
   "단순 메모리가 아니라, 추적 가능한 '믿음'의 지식 베이스", sz=13, c=RGBColor(0x88, 0xAA, 0xCC), fn='Georgia', i=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 17 — WHAT IS MEMORY WIKI + BELIEF LAYER
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
hdr(slide, "Memory Wiki란? — Belief Layer의 구현")
ftr(slide, 17)

# Left: Problem vs Solution
add_rect(slide, Inches(0.6), Inches(1.2), Inches(5.8), Inches(1.8), OFF_WHITE, LIGHT_GREY)
tb(slide, Inches(0.8), Inches(1.3), Inches(5.4), Inches(0.3),
   "기존 메모리의 한계", sz=14, c=SIGNAL_RED, b=True, fn='Georgia')
rule(slide, Inches(0.8), Inches(1.65), Inches(1.5), SIGNAL_RED, Pt(1.5))
rich(slide, Inches(0.8), Inches(1.8), Inches(5.4), Inches(1),
    [("• MEMORY.md는 마크다운 파일 더미 — 뭐가 사실인지 추적 불가", 10, DARK_GREY),
     ("• 확실한 사실과 추측의 구분이 없음", 10, DARK_GREY),
     ("• 정보가 어디서 왔는지 출처를 모름", 10, DARK_GREY),
     ("• 모순되는 정보가 있어도 감지 못함", 10, DARK_GREY)])

add_rect(slide, Inches(6.8), Inches(1.2), Inches(5.8), Inches(1.8), NAVY)
tb(slide, Inches(7.0), Inches(1.3), Inches(5.4), Inches(0.3),
   "Memory Wiki의 해결", sz=14, c=GOLD, b=True, fn='Georgia')
rule(slide, Inches(7.0), Inches(1.65), Inches(1.5), GOLD, Pt(1.5))
rich(slide, Inches(7.0), Inches(1.8), Inches(5.4), Inches(1),
    [("• 각 사실에 출처(evidence) + 확신도(confidence) 부착", 10, WHITE),
     ("• '믿음(belief)' 상태를 명시적으로 추적", 10, WHITE),
     ("• 모순 자동 감지 → contradictions.md 리포트", 10, WHITE),
     ("• 카파시(Karpathy)의 Belief Layer 철학 구현", 10, RGBColor(0xAA,0xCC,0xE8))])

# Core concept
add_rect(slide, Inches(0.6), Inches(3.3), Inches(12.1), Inches(1.2), NAVY)
tb(slide, Inches(0.8), Inches(3.4), Inches(11.7), Inches(0.3),
   "핵심 철학: \"LLM은 flat 로그가 아니라 구조화된 믿음 상태를 유지해야 한다\" — Karpathy",
   sz=13, c=WHITE, b=True, fn='Georgia')
tb(slide, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.5),
   "사실을 텍스트가 아니라 claim + evidence 체인으로 저장. hallucination 감소, 추론 신뢰성 향상, 모순 자동 감지.",
   sz=11, c=RGBColor(0xAA,0xCC,0xE8))

# 2-layer architecture
tb(slide, Inches(0.6), Inches(4.8), Inches(12), Inches(0.3),
   "2-레이어 분리 구조", sz=14, c=NAVY, b=True, fn='Georgia')
rule(slide, Inches(0.6), Inches(5.15), Inches(2), NAVY, Pt(1.5))

# Active Memory box
add_rect(slide, Inches(0.6), Inches(5.4), Inches(5.8), Inches(1.6), OFF_WHITE, TEAL, Pt(1.5))
tb(slide, Inches(0.8), Inches(5.5), Inches(5.4), Inches(0.3),
   "Active Memory Plugin", sz=13, c=TEAL, b=True, fn='Georgia')
tb(slide, Inches(0.8), Inches(5.8), Inches(5.4), Inches(1),
   "리콜, 검색, 드리밍, 프로모션\n\"원료\"를 생산\nmemory-core / QMD / Honcho", sz=10, c=DARK_GREY)

# Arrow
tb(slide, Inches(6.2), Inches(5.9), Inches(0.8), Inches(0.5),
   "→", sz=24, c=DARK_GREY, b=True, a=PP_ALIGN.CENTER)

# Wiki box
add_rect(slide, Inches(6.8), Inches(5.4), Inches(5.8), Inches(1.6), OFF_WHITE, PURPLE, Pt(1.5))
tb(slide, Inches(7.0), Inches(5.5), Inches(5.4), Inches(0.3),
   "Memory Wiki", sz=13, c=PURPLE, b=True, fn='Georgia')
tb(slide, Inches(7.0), Inches(5.8), Inches(5.4), Inches(1),
   "출처, 확신도, 모순 추적\n\"정제된 제품\"\nentities / concepts / syntheses / claims", sz=10, c=DARK_GREY)


# ═══════════════════════════════════════════════════════════
# SLIDE 18 — CORE CONCEPTS: INGEST / ENTITY / CONCEPT / SYNTHESIS / CLAIM
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
hdr(slide, "핵심 개념: 5가지 구성 요소")
ftr(slide, 18)

concepts = [
    ("Source", "수집 (Ingest)", TEAL, [
        "외부 자료를 wiki로 가져오는 과정",
        "파일 또는 URL에서 ingest",
        "sources/ 폴더에 원본 보존",
        "bridge 모드: QMD에서 자동 수집",
    ]),
    ("Entity", "엔티티", NAVY, [
        "추적 가능한 '사물'",
        "person / team / system / project",
        "고유 ID + alias + 연락처",
        "bestUsedFor / notEnoughFor 라우팅",
    ]),
    ("Concept", "컨셉트", SOFT_BLUE, [
        "아이디어 / 패턴 / 정책",
        "entity = 존재, concept = 앎",
        "정의를 명확히 추적",
        "다른 concept과 관계 맺음",
    ]),
    ("Synthesis", "신세시스", AMBER, [
        "여러 소스의 컴파일된 종합 요약",
        "sourceIds로 출처 전체 추적",
        "흩어진 정보를 한 곳에",
        "일반 노트와 근본적 차이",
    ]),
    ("Claim", "클레임", SIGNAL_RED, [
        "\"나는 X를 믿는다\"는 명시적 선언",
        "status: supported / contested /",
        "  superseded / open",
        "evidence + confidence 필수",
    ]),
]

for i, (name, korean, color, bullets) in enumerate(concepts):
    x = Inches(0.4) + i * Inches(2.55)
    y = Inches(1.2)
    
    add_rect(slide, x, y, Inches(2.35), Inches(5.8), OFF_WHITE, LIGHT_GREY, Pt(0.5))
    add_rect(slide, x, y, Inches(2.35), Inches(0.06), color)
    
    # Name
    tb(slide, x+Inches(0.15), y+Inches(0.2), Inches(2.05), Inches(0.4),
       name, sz=16, c=color, b=True, fn='Georgia')
    tb(slide, x+Inches(0.15), y+Inches(0.6), Inches(2.05), Inches(0.25),
       korean, sz=9, c=MED_GREY, b=True)
    rule(slide, x+Inches(0.15), y+Inches(0.9), Inches(1), color, Pt(1.5))
    
    for j, bullet in enumerate(bullets):
        tb(slide, x+Inches(0.15), y+Inches(1.1)+j*Inches(0.55), Inches(2.05), Inches(0.5),
           f"• {bullet}", sz=8, c=DARK_GREY)

# Bottom flow
add_rect(slide, Inches(0.4), Inches(7.1)-Inches(0.7), Inches(12.3), Inches(0.5), NAVY)
tb(slide, Inches(0.6), Inches(7.1)-Inches(0.65), Inches(12), Inches(0.4),
   "흐름:  Source (수집)  →  Entity / Concept (분류)  →  Claim (믿음 선언)  →  Synthesis (종합)  →  Compile (컴파일)",
   sz=11, c=WHITE, b=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 19 — CLAIMS DEEP DIVE + VAULT MODES + OBSIDIAN
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
hdr(slide, "Structured Claims + Vault 모드 + Obsidian 연동")
ftr(slide, 19)

# Left: Claims
add_rect(slide, Inches(0.6), Inches(1.2), Inches(6), Inches(2.8), OFF_WHITE, LIGHT_GREY, Pt(0.5))
tb(slide, Inches(0.8), Inches(1.3), Inches(5.6), Inches(0.3),
   "Claim 상태 라이프사이클", sz=13, c=SIGNAL_RED, b=True, fn='Georgia')
rule(slide, Inches(0.8), Inches(1.65), Inches(1.5), SIGNAL_RED, Pt(1.5))

rich(slide, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2),
    [("supported — 현재 믿음 (확실한 사실)", 10, GREEN_ACC, True),
     ("contested — 논란 있음 (반대 증거 발견)", 10, AMBER, True),
     ("superseded — 대체됨 (새 정보로 교체)", 10, MED_GREY, True),
     ("open — 미확인 (아직 판단 안 됨)", 10, SOFT_BLUE, True),
     ("", 6, DARK_GREY),
     ("모순 감지: 같은 주제에 supported가 2개 이상 → contradictions.md 자동 집계", 9, SIGNAL_RED),
     ("해소 방법: 하나를 superseded로 명시적 처리", 9, DARK_GREY)], fn='Calibri')

# Right: Vault modes
add_rect(slide, Inches(6.8), Inches(1.2), Inches(5.8), Inches(2.8), OFF_WHITE, LIGHT_GREY, Pt(0.5))
tb(slide, Inches(7.0), Inches(1.3), Inches(5.4), Inches(0.3),
   "Vault 모드 3종", sz=13, c=NAVY, b=True, fn='Georgia')
rule(slide, Inches(7.0), Inches(1.65), Inches(1.5), NAVY, Pt(1.5))

rich(slide, Inches(7.0), Inches(1.8), Inches(5.4), Inches(2),
    [("isolated (권장 시작점)", 11, TEAL, True),
     ("독립 vault, 자체 소스, plugin 의존 없음", 9, DARK_GREY),
     ("", 6, DARK_GREY),
     ("bridge (QMD 연동)", 11, SOFT_BLUE, True),
     ("active memory plugin의 공개 산출물을 wiki로 가져옴", 9, DARK_GREY),
     ("QMD public artifacts → sources/ → compile", 9, DARK_GREY),
     ("", 6, DARK_GREY),
     ("unsafe-local (실험적)", 11, SIGNAL_RED, True),
     ("로컬 파일시스템 직접 접근 — 신뢰 경계 파괴 위험", 9, DARK_GREY)], fn='Calibri')

# Bottom: Obsidian
add_rect(slide, Inches(0.6), Inches(4.2), Inches(12.1), Inches(3), NAVY)
tb(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(0.3),
   "Obsidian 연동", sz=14, c=GOLD, b=True, fn='Georgia')
rule(slide, Inches(0.8), Inches(4.65), Inches(1.5), GOLD, Pt(1.5))

# Two patterns
rich(slide, Inches(0.8), Inches(4.8), Inches(5.4), Inches(2.2),
    [("패턴 A (추천): 기존 볼트 내 하위 폴더", 11, RGBColor(0x88,0xDD,0x88), True),
     ("vault.path를 기존 Obsidian 볼트 내 wiki/로 지정", 9, WHITE),
     ("→ Obsidian에서 바로 열리고 편집 가능", 9, WHITE),
     ("", 6, WHITE),
     ("패턴 B: 별도 wiki vault", 11, RGBColor(0xAA,0xCC,0xE8), True),
     ("vault.path를 ~/.openclaw/wiki/main으로 설정", 9, WHITE),
     ("→ Obsidian에 새 볼트로 추가 (vaultName 지정)", 9, WHITE)], fn='Calibri')

# Right: Obsidian features
rich(slide, Inches(6.8), Inches(4.8), Inches(5.4), Inches(2.2),
    [("renderMode: \"obsidian\" 효과:", 11, GOLD, True),
     ("• [[wikilinks]] 형식 백링크 생성", 9, WHITE),
     ("• Obsidian Properties 호환 YAML frontmatter", 9, WHITE),
     ("", 6, WHITE),
     ("useOfficialCli: true 시 추가 기능:", 11, GOLD, True),
     ("• obsidian search / open / daily 명령 사용 가능", 9, WHITE),
     ("• Obsidian 앱 내에서 직접 검색, 페이지 열기", 9, WHITE),
     ("• openclaw wiki obsidian status 로 가용 확인", 9, WHITE)], fn='Calibri')


# ═══════════════════════════════════════════════════════════
# SLIDE 20 — SETUP + SEARCH + DASHBOARDS + WORKFLOW
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
hdr(slide, "설정 / 검색 / 대시보드 / 워크플로우")
ftr(slide, 20)

# Top-left: Setup
add_rect(slide, Inches(0.6), Inches(1.2), Inches(4), Inches(2.5), OFF_WHITE, LIGHT_GREY, Pt(0.5))
tb(slide, Inches(0.8), Inches(1.3), Inches(3.6), Inches(0.3),
   "초기 설정 (5단계)", sz=13, c=NAVY, b=True, fn='Georgia')
rule(slide, Inches(0.8), Inches(1.65), Inches(1.2), NAVY, Pt(1.5))
rich(slide, Inches(0.8), Inches(1.8), Inches(3.6), Inches(1.8),
    [("① openclaw wiki doctor", 10, DARK_GREY, False, PP_ALIGN.LEFT, False),
     ("② openclaw wiki init", 10, DARK_GREY, False, PP_ALIGN.LEFT, False),
     ("③ openclaw wiki ingest ./notes/", 10, DARK_GREY, False, PP_ALIGN.LEFT, False),
     ("④ openclaw wiki compile", 10, DARK_GREY, False, PP_ALIGN.LEFT, False),
     ("⑤ openclaw wiki lint", 10, DARK_GREY, False, PP_ALIGN.LEFT, False)],
    fn='Consolas')

# Top-center: Search modes
add_rect(slide, Inches(4.9), Inches(1.2), Inches(4), Inches(2.5), OFF_WHITE, LIGHT_GREY, Pt(0.5))
tb(slide, Inches(5.1), Inches(1.3), Inches(3.6), Inches(0.3),
   "wiki_search 5가지 모드", sz=13, c=TEAL, b=True, fn='Georgia')
rule(slide, Inches(5.1), Inches(1.65), Inches(1.2), TEAL, Pt(1.5))
rich(slide, Inches(5.1), Inches(1.8), Inches(3.6), Inches(1.8),
    [("auto — 기본 (모르면 이걸)", 9, DARK_GREY),
     ("find-person — 사람/핸들/이메일", 9, DARK_GREY),
     ("route-question — \"누구한테 물어봐?\"", 9, DARK_GREY),
     ("source-evidence — 출처/근거 추적", 9, DARK_GREY),
     ("raw-claim — 클레임 직접 조회", 9, DARK_GREY)],
    fn='Calibri')

# Top-right: Dashboards
add_rect(slide, Inches(9.2), Inches(1.2), Inches(3.5), Inches(2.5), OFF_WHITE, LIGHT_GREY, Pt(0.5))
tb(slide, Inches(9.4), Inches(1.3), Inches(3.1), Inches(0.3),
   "대시보드 9종", sz=13, c=PURPLE, b=True, fn='Georgia')
rule(slide, Inches(9.4), Inches(1.65), Inches(1.2), PURPLE, Pt(1.5))
rich(slide, Inches(9.4), Inches(1.8), Inches(3.1), Inches(1.8),
    [("contradictions — 모순 클러스터", 8, SIGNAL_RED, True),
     ("claim-health — 근거 없는 믿음", 8, AMBER, True),
     ("person-directory — 사람 라우팅", 8, TEAL, True),
     ("low-confidence — 저신뢰도", 8, DARK_GREY),
     ("stale-pages — 낡은 정보", 8, DARK_GREY),
     ("open-questions — 미해결 질문", 8, DARK_GREY),
     ("relationship-graph — 관계", 8, DARK_GREY),
     ("provenance-coverage — 출처", 8, DARK_GREY),
     ("privacy-review — 프라이버시", 8, DARK_GREY)],
    fn='Calibri')

# Bottom: Workflow
add_rect(slide, Inches(0.6), Inches(4.0), Inches(12.1), Inches(1.5), NAVY)
tb(slide, Inches(0.8), Inches(4.1), Inches(11.7), Inches(0.3),
   "추천 워크플로우", sz=14, c=GOLD, b=True, fn='Georgia')
rule(slide, Inches(0.8), Inches(4.45), Inches(1.5), GOLD, Pt(1.5))

steps = [
    ("① isolated로 시작", "vault 안정화, 기본 구조 익히기"),
    ("② Obsidian 연동", "renderMode: obsidian, 기존 볼트 연결"),
    ("③ bridge로 전환", "QMD + wiki 조합, 자동 수집"),
    ("④ 대시보드 활용", "lint 정기 실행, 모순/신뢰도 모니터링"),
]
for i, (step, desc) in enumerate(steps):
    x = Inches(0.8) + i * Inches(3)
    tb(slide, x, Inches(4.6), Inches(2.8), Inches(0.25),
       step, sz=11, c=RGBColor(0x88,0xDD,0x88), b=True)
    tb(slide, x, Inches(4.85), Inches(2.8), Inches(0.25),
       desc, sz=9, c=RGBColor(0xAA,0xCC,0xE8))

# Bottom: Lint
add_rect(slide, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.5), OFF_WHITE, LIGHT_GREY, Pt(0.5))
tb(slide, Inches(0.8), Inches(5.8), Inches(3), Inches(0.3),
   "wiki_lint 검사 항목", sz=13, c=NAVY, b=True, fn='Georgia')
rule(slide, Inches(0.8), Inches(6.15), Inches(1.5), NAVY, Pt(1.5))

lint_items = [
    ("구조", "vault 레이아웃, 필수 파일 누락"),
    ("출처 갭", "claim에 evidence 없음"),
    ("모순", "동일 주제 supported 2개 이상"),
    ("미해소 질문", "open questions 방치"),
    ("저신뢰도", "confidence < 0.5"),
    ("낡은 정보", "lastRefreshedAt 오래됨"),
    ("프라이버시", "sensitive 항목 노출 위험"),
]

for i, (cat, desc) in enumerate(lint_items):
    col = i % 4
    row = i // 4
    x = Inches(0.8) + col * Inches(3)
    y = Inches(6.3) + row * Inches(0.45)
    tb(slide, x, y, Inches(0.8), Inches(0.2), cat, sz=8, c=NAVY, b=True)
    tb(slide, x+Inches(0.85), y, Inches(2), Inches(0.2), desc, sz=8, c=DARK_GREY)


# ── Update footer page numbers for all slides ──
# Re-number all slides to reflect 20 total
# (We just added slides 16-20, existing 1-15 already have correct numbers)

# ── Save ──
output = "/Users/conanssam-m4/Downloads/openclaw-memory-mckinsey.pptx"
prs.save(output)
print(f"✅ Saved: {output} ({len(prs.slides)} slides)")
