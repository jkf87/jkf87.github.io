#!/usr/bin/env python3
"""Add speaker notes (Korean script) to existing McKinsey PPTX"""

from pptx import Presentation
from pptx.util import Pt

prs = Presentation('/Users/conanssam-m4/Downloads/openclaw-memory-mckinsey.pptx')

scripts = [
    # Slide 1: Title
    """안녕하세요. 오늘은 OpenClaw의 메모리 시스템을 정리해보려고 합니다.

주제는 하나예요. AI 에이전트가 어떻게 기억하는가.

지금까지 에이전트 쓰면서 가장 답답했던 게 뭔가요? 어제 말한 걸 오늘 또 설명해야 하는 거죠. 이 문제를 OpenClaw가 5개 레이어로 어떻게 풀었는지, 오늘 같이 보겠습니다.""",

    # Slide 2: Problem
    """먼저 문제부터 명확히 하겠습니다.

첫째, 에이전트는 매 세션이 끝나면 모든 걸 잊어버립니다. 대화 창을 닫으면 그게 끝이에요.

둘째, 기억을 저장하려면 우리가 직접 "기억해"라고 말해야 합니다. 에이전트가 스스로 판단해서 저장하는 게 아니에요.

셋째, 파일에 기억이 있어도 찾지 못하면 의미가 없습니다. "인증"이라고 검색했는데 JWT 설정이 안 나오면 그건 검색이 깨진 거죠.

넷째, 전체 컨텍스트를 매번 복사하면 토큰이 터집니다. 연간 1,950만 토큰. 비용도 문제지만 윈도우 한계 때문에 아예 불가능합니다.

이게 우리가 풀어야 할 문제입니다.""",

    # Slide 3: 5 Layer Architecture
    """그래서 OpenClaw는 5개 레이어로 이 문제에 접근합니다.

L1은 파일 기반 메모리입니다. MEMORY.md와 일일 노트. 가장 기본이 되는 층이에요.

L2는 검색 엔진입니다. BM25 키워드 검색과 Vector 의미 검색을 합친 하이브리드 방식이고, 4종 백엔드 중 선택할 수 있습니다.

L3은 Active Memory입니다. 사용자가 질문하기도 전에 에이전트가 먼저 기억을 찾아서 주입하는 겁니다.

L4는 Dreaming입니다. 사람이 자면서 기억을 정리하듯이, 백그라운드에서 단기 기억을 장기 기억으로 승격시킵니다.

L5는 Context Engine입니다. 모델이 볼 컨텍스트를 지능적으로 조립하고, 대화가 길어지면 압축합니다.

각 레이어가 독립적으로 작동해서, 일부만 켜도 되고 전부 켜도 됩니다.""",

    # Slide 4: L1 File-Based
    """L1부터 보겠습니다. 파일 기반 메모리.

MEMORY.md는 장기 기억입니다. 중요한 결정, 선호, 사실을 큐레이션해서 저장해요. 매 DM 세션이 시작될 때 자동으로 로딩됩니다.

일일 노트는 memory 폴더 아래 날짜별로 저장됩니다. 오늘과 어제 파일이 자동 로딩되고, Dreaming의 원료로도 쓰입니다.

DREAMS.md는 Dreaming 시스템이 출력하는 휴먼 리뷰용 파일입니다. 뭐가 승격되었는지 사람이 확인할 수 있어요.

핵심은 이겁니다. 마크다운 파일 자체가 메모리입니다. DB가 아니에요. 언제든 직접 열어서 수정할 수 있고, 에이전트도 직접 편집합니다.""",

    # Slide 5: L2 Search 4 Backends
    """L2, 검색 엔진입니다. 4종 백엔드가 있습니다.

Builtin이 기본값입니다. SQLite에 BM25 키워드 검색과 Vector 의미 검색을 합친 하이브리드 방식이에요. API 키 하나만 있으면 자동으로 활성화됩니다. 대부분의 사용자가 여기서 만족합니다.

QMD는 Builtin 위에 리랭킹과 쿼리 확장을 추가합니다. 외부 디렉토리 인덱싱도 되고, 과거 대화 검색도 됩니다. 완전 로컬이라 API 키가 필요 없어요.

Honcho는 AI가 사용자를 자동으로 모델링합니다. "기억해" 안 해도 알아서 이해하고, 채널 경계 없이 기억이 따라옵니다.

LanceDB는 자동 리콜과 자동 캡처가 핵심입니다. 매 턴마다 알아서 기억을 찾고 저장합니다. S3에 저장하면 여러 기기에서도 동기화됩니다.""",

    # Slide 6: Hybrid Search
    """하이브리드 검색이 어떻게 작동하는지 보겠습니다.

사용자가 질문하면, 그게 두 갈래로 나뉩니다.

BM25는 키워드 매칭입니다. 정확한 함수명이나 에러 문자열을 찾을 때 강해요. "auth_middleware"라고 검색하면 그 파일을 정확히 찾습니다.

Vector는 의미 검색입니다. "데이터베이스 성능"이라고 검색하면 N+1 쿼리 수정 기록을 찾아줍니다. 단어가 전혀 달라도 의미가 같으면 매칭되는 거죠.

이 두 결과를 RRF, Reciprocal Rank Fusion으로 병합합니다. 서로 다른 강점을 합치는 거예요.

CJK 트라이그램 덕분에 한국어도 잘 됩니다. 8개 임베딩 프로바이더를 지원하고, API 키만 있으면 자동 감지합니다.""",

    # Slide 7: QMD
    """QMD를 좀 더 깊이 보겠습니다. 토비가 만든 로컬 검색 사이드카입니다.

Builtin 위에 추가되는 게 세 가지 있습니다.

리랭킹은 1차 검색 후 별도 모델로 재평가하는 겁니다. 정확도가 높은 결과가 위로 올라와요.

쿼리 확장은 "인증"이라고 검색하면 JWT, OAuth, 미들웨어까지 자동으로 연관 쿼리를 만들어줍니다.

그리고 워크스페이스 밖에 있는 ~/notes 같은 디렉토리도 인덱싱할 수 있고, 과거 대화 기록도 검색 가능합니다.

장애 대응도 중요합니다. QMD가 죽으면 자동으로 Builtin으로 폴백합니다. 에이전트는 계속 정상 작동해요.""",

    # Slide 8: Honcho
    """Honcho는 근본적으로 다른 접근입니다.

다른 백엔드는 텍스트를 청크로 나누고 인덱싱하는 "검색 중심"이에요. Honcho는 대화 전체를 관찰해서 사용자를 모델링하는 "이해 중심"입니다.

매 턴마다 사용자 메시지와 에이전트 메시지를 관찰로 저장합니다. 그러면 Honcho가 "이 사람은 TypeScript를 선호하는구나"를 자동으로 학습합니다.

멀티 에이전트 인식도 됩니다. 부모 에이전트가 자식 에이전트를 자동 추적하고, Telegram에서 Discord로 채널을 바꿔도 같은 사용자 컨텍스트가 유지됩니다.

"기억해"라고 안 해도 된다는 게 핵심이에요. 에이전트가 알아서 사용자를 이해합니다.""",

    # Slide 9: LanceDB
    """LanceDB의 핵심은 autoRecall과 autoCapture입니다.

autoRecall은 매 턴 전에 사용자 메시지를 쿼리로 임베딩해서, LanceDB에서 유사한 기억을 찾아 자동으로 컨텍스트에 주입합니다.

autoCapture는 매 턴 후에 어시스턴트 응답이 기준 이하면 자동으로 저장합니다. 사용자가 "기억해"라고 안 해도 돼요.

Ollama 임베딩을 네이티브로 지원합니다. mxbai-embed-large 같은 로컬 모델을 쓰면 완전 무료예요.

그리고 S3에 저장할 수 있습니다. dbPath를 s3://로 설정하면 여러 기기에서 같은 메모리를 공유합니다.

Builtin과의 차이를 명확히 말씀드리면, Builtin은 수동 검색이고 LanceDB는 자동입니다. 설정만 하면 알아서 작동해요.""",

    # Slide 10: Active Memory
    """L3, Active Memory입니다. 이게 제가 가장 좋아하는 기능이에요.

보통은 에이전트가 memory_search를 직접 호출해야 기억을 찾습니다. 근데 Active Memory는 사용자가 질문하기도 전에, 별도 서브에이전트가 먼저 기억을 찾습니다.

Step 1, 사용자가 질문합니다.
Step 2, 차단 서브에이전트가 메모리를 검색합니다.
Step 3, 관련 기억이 있으면 숨겨서 프롬프트에 주입합니다.
Step 4, 메인 에이전트가 자연스럽게 기억을 활용해서 답변합니다.

쿼리 모드는 세 가지가 있습니다. message는 마지막 메시지만 봐서 빠르고, recent는 최근 대화 꼬리까지 봐서 균형 잡혀 있고, full은 전체 대화를 봐서 가장 정확합니다.

빠른 모델 권장합니다. gemini-flash나 cerebras 같은 걸 쓰면 지연이 거의 안 느껴져요.""",

    # Slide 11: Dreaming
    """L4, Dreaming입니다. 인간의 수면 기억 정리를 모방합니다.

Light Sleep은 최근 단기 기억을 정리하고 승격 후보를 스테이징합니다. MEMORY.md에는 안 써요.

Deep Sleep이 핵심입니다. 6개 시그널로 점수를 매겨서 장기 기억으로 승격시킵니다. 관련성 30%, 빈도 24%, 쿼리 다양성 15%, 최신성 15%, 통합 강도 10%, 개념 풍부도 6%.

그리고 3개 게이트를 다 통과해야 승격됩니다. 최소 점수, 최소 리콜 횟수, 최소 고유 쿼리 수. 아무거나 다 저장하는 게 아니에요.

REM Sleep은 패턴을 추출하고 Deep 랭킹에 피드백을 줍니다.

기본 비활성화입니다. 매일 새벽 3시에 자동 실행되고, DREAMS.md에서 사람이 리뷰할 수 있어요.""",

    # Slide 12: Context Engine
    """마지막 L5, Context Engine과 Compaction입니다.

Context Engine은 4개 훅으로 작동합니다. ingest는 메시지를 저장하고, assemble은 토큰 예산 안에서 모델이 볼 메시지를 조립합니다. compact는 대화가 길어지면 요약하고, afterTurn은 실행 후 상태를 업데이트합니다.

Compaction 프로세스가 중요합니다. 컨텍스트 한계에 가까워지면, 먼저 에이전트에게 "중요한 거 파일에 저장해"라고 자동으로提醒합니다. 그 다음 오래된 대화를 요약으로 압축하고, 최근 메시지는 그대로 보존합니다. 전체 기록은 디스크에 그대로 남아서 나중에 검색할 수 있어요.

플러그인으로 교체도 가능합니다. 기본은 legacy 엔진이고, lossless-claw 같은 서드파티 엔진으로 바꿀 수도 있어요.""",

    # Slide 13: Comparison
    """4개 백엔드를 한 표로 비교하겠습니다.

검색은 Builtin과 QMD가 BM25+Vector, Honcho는 의미 검색, LanceDB는 Vector만 사용합니다.

자동 리콜은 Honcho와 LanceDB만 지원합니다. 자동 캡처도 마찬가지입니다.

사용자 모델링은 Honcho만 제공합니다. 에이전트가 사용자를 자동으로 이해하는 기능이죠.

과거 대화 검색은 QMD와 Honcho, 외부 디렉토리는 QMD만 가능합니다.

클라우드 저장은 Honcho가 매니지드 API, LanceDB가 S3를 지원합니다.

결론: Builtin만으로도 충분하고, 필요에 따라 하나씩 추가하면 됩니다.""",

    # Slide 14: Practical Guide
    """실전 가이드입니다. 세 단계로 나눴습니다.

기본은 0분이에요. API 키 하나만 설정하면 Builtin 하이브리드 검색이 자동 활성화됩니다. MEMORY.md와 일일 노트만으로도 충분합니다. 대부분의 사용자가 여기서 만족합니다.

중급은 10분입니다. Active Memory를 추가하세요. "기억해" 안 해도 에이전트가 자동으로 기억을 활용합니다. 빠른 모델 하나만 설정하면 돼요.

고급은 30분입니다. QMD로 완전 로컬 검색, Honcho로 사용자 모델링, LanceDB로 자동 리콜/캡처, 그리고 Dreaming까지. 전부 켜면 인간처럼 기억합니다.

핵심은 점진적 업그레이드입니다. 기본에서 시작해서 필요할 때 하나씩 추가하세요.""",

    # Slide 15: Closing
    """마지막입니다.

기억하는 에이전트와 기억하지 않는 에이전트. 이 차이가 사용자 경험의 전부입니다.

같은 걸 10번 설명해야 하는 에이전트와, 한 번 말하면 알아서 기억하고 다음에 자연스럽게 활용하는 에이전트. 어떤 게 좋을지는 자명하죠.

OpenClaw는 5개 레이어로 이 차이를 메웁니다. 기본 설정만으로도 강력하고, 모두 켜면 인간처럼 기억합니다.

감사합니다.""",

    # Slide 16: Wiki Section Divider
    """여기서부터 두 번째 파트입니다. Memory Wiki.

앞에서 본 5개 레이어가 "기억하는" 시스템이었다면, Memory Wiki는 그 위에 얹는 "아는" 시스템입니다.

단순히 기억하는 걸 넘어서, 뭐가 사실인지, 얼마나 확실한지, 어디서 나온 정보인지를 추적하는 지식 베이스입니다.

카파시의 Belief Layer 철학을 실제로 구현한 거예요.""",

    # Slide 17: What is Memory Wiki
    """Memory Wiki가 뭔지 명확히 하겠습니다.

기존 MEMORY.md의 한계는 이겁니다. 마크다운 파일 더미인데, 뭐가 사실인지, 확실한 건지, 어디서 들은 건지 전혀 추적이 안 돼요.

Memory Wiki는 각 사실에 출처와 확신도를 붙입니다. "믿음" 상태를 명시적으로 추적하는 거죠.

2-레이어 분리가 핵심입니다. Active Memory Plugin은 리콜, 검색, 드리밍을 담당하고, Memory Wiki는 출처, 확신도, 모순 추적을 담당합니다. 서로 대체하는 게 아니라 옆에 앉는 겁니다.

wiki가 없어도 에이전트는 잘 작동합니다. 근데 wiki가 있으면 에이전트가 "이건 확실해"와 "이건 추측이야"를 구분할 수 있게 됩니다.""",

    # Slide 18: 5 Core Concepts
    """5가지 핵심 구성 요소입니다.

Source는 원본 자료입니다. 파일이나 URL에서 수집해서 sources 폴더에 저장합니다. bridge 모드면 QMD에서 자동으로 들어와요.

Entity는 추적 대상입니다. 사람, 팀, 시스템, 프로젝트. 각각 고유 ID와 alias, 연락처를 가집니다.

Concept은 아이디어와 패턴입니다. Entity가 "존재하는 것"이라면, Concept은 "알고 있는 것"입니다.

Synthesis는 여러 소스에서 컴파일된 종합 요약입니다. 흩어진 정보를 한 곳에 모으는 거죠.

Claim이 가장 중요합니다. "나는 X를 믿는다"는 명시적 선언입니다. 각 claim에는 출처와 확신도가 필수로 붙습니다.

흐름은 Source에서 시작해서 Entity와 Concept으로 분류되고, Claim으로 믿음을 선언하고, Synthesis로 종합됩니다.""",

    # Slide 19: Claims + Vault + Obsidian
    """Claim의 상태 라이프사이클을 보겠습니다.

supported는 현재 믿는 사실입니다. contested는 반대 증거가 나온 상태이고요. superseded는 새 정보로 대체된 거예요. open은 아직 판단 안 된 겁니다.

모순 감지가 어떻게 되냐면, 같은 주제에 supported claim이 2개 이상 있으면 contradictions.md에 자동 집계됩니다. 하나를 superseded로 명시적으로 처리해야 해요. 이게 단순 로그와 근본적으로 다른 점입니다.

Vault 모드는 세 가지입니다. isolated가 권장 시작점이고, bridge는 QMD 연동할 때 쓰고, unsafe-local은 위험하니 특수한 경우만 쓰세요.

Obsidian 연동은 두 패턴이 있습니다. 기존 볼트 내에 wiki 폴더를 만드는 걸 추천합니다. renderMode를 obsidian으로 설정하면 wikilinks도 되고 Obsidian Properties도 호환돼요.""",

    # Slide 20: Setup + Search + Dashboards
    """마지막 슬라이드입니다. 실전 설정부터 활용까지.

초기 설정은 5단계입니다. doctor로 진단하고, init으로 만들고, ingest로 자료를 넣고, compile하고, lint로 검증합니다.

검색은 5가지 모드가 있습니다. auto가 기본이고, 사람 찾을 땐 find-person, "누구한테 물어봐?"는 route-question, 출처 추적은 source-evidence, 클레임 직접 조회는 raw-claim입니다.

대시보드는 9종이 자동 생성됩니다. 모순, 클레임 건강, 사람 디렉토리가 top 3입니다.

추천 워크플로우: isolated로 시작하고, Obsidian 연동하고, 안정되면 bridge로 전환하고, 대시보드로 모니터링하세요.

이상으로 발표를 마치겠습니다. 감사합니다.""",
]

for i, slide in enumerate(prs.slides):
    if i >= len(scripts):
        break
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = scripts[i]
    p.font.size = Pt(12)

output = "/Users/conanssam-m4/Downloads/openclaw-memory-mckinsey.pptx"
prs.save(output)
print(f"✅ Notes added to {len(prs.slides)} slides → {output}")
