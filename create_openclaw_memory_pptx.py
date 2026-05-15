#!/usr/bin/env python3
"""OpenClaw Memory System PPTX — Glassmorphism style"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── Constants ──────────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

# Glassmorphism palette
BG_DEEP    = RGBColor(0x0F, 0x0F, 0x2D)
GLASS_FILL = RGBColor(0xFF, 0xFF, 0xFF)  # will set alpha separately
GLASS_BORDER = RGBColor(0xFF, 0xFF, 0xFF)
TITLE_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BODY_WHITE   = RGBColor(0xE0, 0xE0, 0xF0)
ACCENT_CYAN  = RGBColor(0x67, 0xE8, 0xF9)
ACCENT_VIOLET = RGBColor(0xA7, 0x8B, 0xFA)
ACCENT_GREEN  = RGBColor(0x34, 0xD3, 0x99)
ACCENT_AMBER  = RGBColor(0xFB, 0xBF, 0x24)
ACCENT_ROSE   = RGBColor(0xFB, 0x71, 0x85)
GLOW_CYAN     = RGBColor(0x22, 0xD3, 0xEE)
GLOW_VIOLET   = RGBColor(0x8B, 0x5C, 0xF6)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# Remove default layouts — use blank
blank_layout = prs.slide_layouts[6]  # blank

def set_slide_bg(slide, color=BG_DEEP):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def set_shape_alpha(shape, alpha_pct):
    """Set fill alpha (0-100) on a shape via XML."""
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    sp = shape._element
    srgb = sp.find('.//' + qn('a:srgbClr'))
    if srgb is not None:
        # Remove existing alpha
        for old in srgb.findall(qn('a:alpha')):
            srgb.remove(old)
        alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
        alpha_elem.set('val', str(int(alpha_pct * 1000)))

def add_glow_circle(slide, left, top, width, height, color, alpha=15):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    set_shape_alpha(shape, alpha)
    return shape

def add_glass_card(slide, left, top, width, height, alpha=18):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape.line.width = Pt(1.5)
    # Set fill alpha
    set_shape_alpha(shape, alpha)
    # Set line alpha
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    ln = shape._element.find('.//' + qn('a:ln'))
    if ln is not None:
        sr = ln.find('.//' + qn('a:srgbClr'))
        if sr is not None:
            la = etree.SubElement(sr, qn('a:alpha'))
            la.set('val', str(25 * 1000))
    # Adjust corner radius
    try:
        shape.adjustments[0] = 0.05
    except:
        pass
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=BODY_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_text(slide, left, top, width, height, items, font_size=13,
                    color=BODY_WHITE, spacing=Pt(6), font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_CYAN):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_kpi(slide, left, top, number, label, color=ACCENT_CYAN):
    add_text_box(slide, left, top, Inches(2.5), Inches(0.8), number,
                 font_size=52, color=color, bold=True, alignment=PP_ALIGN.CENTER,
                 font_name='Segoe UI Light')
    add_text_box(slide, left, top + Inches(0.7), Inches(2.5), Inches(0.5), label,
                 font_size=12, color=BODY_WHITE, bold=False, alignment=PP_ALIGN.CENTER)

def add_icon_circle(slide, left, top, size, color, icon_text):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Add text inside
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(int(size / Inches(1) * 16))
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    return shape


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

# Glow blobs
add_glow_circle(slide, Inches(-1), Inches(-1), Inches(6), Inches(6), GLOW_VIOLET, alpha=12)
add_glow_circle(slide, Inches(8), Inches(3), Inches(5), Inches(5), GLOW_CYAN, alpha=10)

# Glass card (center)
card = add_glass_card(slide, Inches(2.5), Inches(1.5), Inches(8.333), Inches(4.5), alpha=12)

# Title
add_text_box(slide, Inches(3), Inches(2), Inches(7.333), Inches(1.5),
             "OpenClaw Memory System",
             font_size=44, color=TITLE_WHITE, bold=True, alignment=PP_ALIGN.CENTER,
             font_name='Segoe UI Light')

# Accent line
add_accent_line(slide, Inches(5.5), Inches(3.4), Inches(2.333), ACCENT_CYAN)

# Subtitle
add_text_box(slide, Inches(3), Inches(3.8), Inches(7.333), Inches(1),
             "AI 에이전트가 기억하는 5개 레이어 아키텍처",
             font_size=20, color=BODY_WHITE, bold=False, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(3), Inches(4.7), Inches(7.333), Inches(0.6),
             "파일 → 검색 → 자동 리콜 → 백그라운드 통합 → 컨텍스트 조립",
             font_size=13, color=ACCENT_CYAN, bold=False, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_glow_circle(slide, Inches(9), Inches(-1), Inches(5), Inches(5), GLOW_VIOLET, alpha=8)

# Section header
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.5),
             "THE PROBLEM", font_size=11, color=ACCENT_CYAN, bold=True,
             font_name='Segoe UI')
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1.5), ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.8),
             "에이전트는 매 세션마다 기억을 잃는다",
             font_size=32, color=TITLE_WHITE, bold=True, font_name='Segoe UI Light')

# Problem cards
problems = [
    ("🔄", "세션 리셋", "새 대화 시작하면 이전 컨텍스트 전부 사라짐.\nCLAUDE.md 200줄 한계로 담을 수 있는 정보가 제한적"),
    ("📋", "수동 기억", "\"기억해\"라고 직접 말해야 파일에 저장.\n에이전트가 스스로 판단해서 기억하지 못함"),
    ("🔍", "검색 누락", "기억이 파일에 있어도 찾지 못하면 의미 없음.\n단순 키워드 검색으로는 맥락을 못 찾음"),
    ("💰", "토큰 낭비", "컨텍스트 전체를 매번 복붙하면 연간 1,950만 토큰.\n비용 폭발 + 윈도우 한계 초과"),
]

for i, (icon, title, desc) in enumerate(problems):
    x = Inches(0.8) + i * Inches(3.1)
    y = Inches(2.5)
    add_glass_card(slide, x, y, Inches(2.8), Inches(4.2), alpha=10)
    add_icon_circle(slide, x + Inches(0.9), y + Inches(0.3), Inches(0.9), ACCENT_ROSE if i == 0 else ACCENT_AMBER if i == 1 else ACCENT_VIOLET if i == 2 else ACCENT_CYAN, icon)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.5), Inches(2.4), Inches(0.5),
                 title, font_size=16, color=TITLE_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(2.2), Inches(2.4), Inches(1.8),
                 desc, font_size=11, color=BODY_WHITE, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — 5 LAYER OVERVIEW
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_glow_circle(slide, Inches(-2), Inches(4), Inches(5), Inches(5), GLOW_CYAN, alpha=8)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.5),
             "ARCHITECTURE", font_size=11, color=ACCENT_VIOLET, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1.5), ACCENT_VIOLET)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.8),
             "5개 레이어 메모리 아키텍처",
             font_size=32, color=TITLE_WHITE, bold=True, font_name='Segoe UI Light')

layers = [
    ("L1", "파일 기반", "MEMORY.md + 일일노트", ACCENT_CYAN),
    ("L2", "검색 엔진", "BM25 + Vector 하이브리드", ACCENT_VIOLET),
    ("L3", "Active Memory", "자동 리콜 (사용자 질문 전)", ACCENT_GREEN),
    ("L4", "Dreaming", "백그라운드 통합 (수면)", ACCENT_AMBER),
    ("L5", "Context Engine", "컨텍스트 조립 + Compaction", ACCENT_ROSE),
]

for i, (num, title, desc, color) in enumerate(layers):
    y = Inches(2.2) + i * Inches(1.0)
    # Layer number circle
    add_icon_circle(slide, Inches(0.8), y, Inches(0.7), color, num)
    # Glass bar
    add_glass_card(slide, Inches(1.8), y, Inches(10.5), Inches(0.8), alpha=8)
    add_text_box(slide, Inches(2.1), y + Inches(0.1), Inches(2.5), Inches(0.6),
                 title, font_size=16, color=TITLE_WHITE, bold=True)
    add_text_box(slide, Inches(5), y + Inches(0.1), Inches(7), Inches(0.6),
                 desc, font_size=13, color=BODY_WHITE)
    # Arrow indicator
    if i < len(layers) - 1:
        add_text_box(slide, Inches(1.05), y + Inches(0.65), Inches(0.5), Inches(0.4),
                     "↓", font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)

# Side note
add_glass_card(slide, Inches(8.5), Inches(2.2), Inches(3.8), Inches(1.8), alpha=10)
add_text_box(slide, Inches(8.7), Inches(2.35), Inches(3.4), Inches(0.4),
             "💡 실전 의미", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(8.7), Inches(2.8), Inches(3.4), Inches(1.1),
             "각 레이어가 독립 작동\n→ 일부만 써도 OK\n→ 모두 켜면 인간처럼 기억\n→ 기본값만으로도 강력",
             font_size=11, color=BODY_WHITE)


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — L1: FILE-BASED MEMORY
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_glow_circle(slide, Inches(8), Inches(4), Inches(5), Inches(5), GLOW_CYAN, alpha=8)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.5),
             "LAYER 1", font_size=11, color=ACCENT_CYAN, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1.5), ACCENT_CYAN)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.8),
             "파일 기반 메모리 — 기본층",
             font_size=30, color=TITLE_WHITE, bold=True, font_name='Segoe UI Light')

# Three file cards
files_data = [
    ("MEMORY.md", "장기 기억", [
        "중요 사실, 결정, 선호",
        "매 DM 세션 시작 시 자동 로딩",
        "\"기억해\" → 자동 저장",
        "에이전트의 큐레이티드 메모리",
    ], ACCENT_CYAN),
    ("memory/YYYY-MM-DD.md", "일일 노트", [
        "당일 관찰, 대화 기록",
        "오늘 + 어제 파일 자동 로딩",
        "세션별 자연 발생 기록",
        "Dreaming의 원료",
    ], ACCENT_VIOLET),
    ("DREAMS.md", "꿈 일기", [
        "Dreaming 시스템 출력",
        "휴먼 리뷰용 (옵션)",
        "승격 후보 + 반영 기록",
        "투명한 메모리 관리",
    ], ACCENT_GREEN),
]

for i, (fname, role, bullets, color) in enumerate(files_data):
    x = Inches(0.8) + i * Inches(4.1)
    y = Inches(2.3)
    add_glass_card(slide, x, y, Inches(3.8), Inches(4.5), alpha=10)
    # File icon + name
    add_icon_circle(slide, x + Inches(0.2), y + Inches(0.3), Inches(0.6), color, "📄")
    add_text_box(slide, x + Inches(1), y + Inches(0.35), Inches(2.5), Inches(0.5),
                 fname, font_size=14, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.1), Inches(3.2), Inches(0.4),
                 role, font_size=16, color=TITLE_WHITE, bold=True)
    add_accent_line(slide, x + Inches(0.3), y + Inches(1.55), Inches(1.5), color)
    add_bullet_text(slide, x + Inches(0.3), y + Inches(1.8), Inches(3.2), Inches(2.5),
                    [f"• {b}" for b in bullets], font_size=11, color=BODY_WHITE)

# Bottom insight
add_glass_card(slide, Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4), alpha=8)
add_text_box(slide, Inches(1), Inches(6.92), Inches(11.3), Inches(0.35),
             "💡 핵심: 마크다운 파일 = 메모리. 복잡한 DB 없이도 작동. 파일을 직접 열고 수정할 수 있다.",
             font_size=11, color=ACCENT_CYAN)


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — L2: SEARCH ENGINE (Overview)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_glow_circle(slide, Inches(-1), Inches(-1), Inches(4), Inches(4), GLOW_VIOLET, alpha=8)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.5),
             "LAYER 2", font_size=11, color=ACCENT_VIOLET, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1.5), ACCENT_VIOLET)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.8),
             "검색 엔진 — 4종 백엔드",
             font_size=30, color=TITLE_WHITE, bold=True, font_name='Segoe UI Light')

# 4 backend cards
backends = [
    ("Builtin", "기본값", "SQLite + FTS5(BM25)\n+ Vector 임베딩\nAPI 키 하나면 자동 활성\nCJK 트라이그램 지원", ACCENT_CYAN),
    ("QMD", "로컬 고급", "Builtin + 리랭킹\n+ 쿼리 확장\n외부 디렉토리 인덱싱\n과거 대화 검색", ACCENT_VIOLET),
    ("Honcho", "AI 네이티브", "자동 사용자 모델링\n멀티 에이전트 인식\n관찰(Observation) 기반\n크로스 세션 기억", ACCENT_GREEN),
    ("LanceDB", "벡터 DB", "LanceDB 컬럼 스토어\n자동 리콜 + 자동 캡처\nOllama 네이티브 지원\nS3 클라우드 저장", ACCENT_AMBER),
]

for i, (name, tag, desc, color) in enumerate(backends):
    x = Inches(0.6) + i * Inches(3.15)
    y = Inches(2.3)
    add_glass_card(slide, x, y, Inches(2.95), Inches(4.5), alpha=10)
    add_icon_circle(slide, x + Inches(0.9), y + Inches(0.3), Inches(0.9), color, name[:1])
    add_text_box(slide, x + Inches(0.2), y + Inches(1.4), Inches(2.55), Inches(0.4),
                 name, font_size=18, color=TITLE_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.8), Inches(2.55), Inches(0.3),
                 tag, font_size=10, color=color, alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, x + Inches(0.6), y + Inches(2.15), Inches(1.5), color)
    add_text_box(slide, x + Inches(0.2), y + Inches(2.4), Inches(2.55), Inches(2),
                 desc, font_size=11, color=BODY_WHITE, alignment=PP_ALIGN.CENTER)

# Bottom insight
add_glass_card(slide, Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4), alpha=8)
add_text_box(slide, Inches(1), Inches(6.92), Inches(11.3), Inches(0.35),
             "💡 핵심: Builtin만으로도 BM25+Vector 하이브리드. QMD/Honcho/LanceDB는 선택적 업그레이드.",
             font_size=11, color=ACCENT_VIOLET)


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — SEARCH ENGINE: How Hybrid Works
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_glow_circle(slide, Inches(7), Inches(0), Inches(5), Inches(5), GLOW_CYAN, alpha=6)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.5),
             "LAYER 2 — DEEP DIVE", font_size=11, color=ACCENT_VIOLET, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1.5), ACCENT_VIOLET)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.8),
             "하이브리드 검색이란?",
             font_size=30, color=TITLE_WHITE, bold=True, font_name='Segoe UI Light')

# Flow diagram
# Query → split → BM25 / Vector → Merge → Results
flow_items = [
    ("질문", Inches(0.5), ACCENT_AMBER),
    ("→", Inches(2.3), BODY_WHITE),
    ("BM25", Inches(2.8), ACCENT_CYAN),
    ("→", Inches(4.8), BODY_WHITE),
    ("RRF", Inches(5.3), ACCENT_VIOLET),
    ("→", Inches(7.1), BODY_WHITE),
    ("결과", Inches(7.6), ACCENT_GREEN),
]

y_flow = Inches(2.5)
for text, x, color in flow_items:
    if text == "→":
        add_text_box(slide, x, y_flow + Inches(0.15), Inches(0.6), Inches(0.5),
                     "→", font_size=28, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_flow, Inches(1.8), Inches(0.7))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.adjustments[0] = 0.15
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = BG_DEEP
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

# Also show vector path
add_text_box(slide, Inches(2.8), y_flow + Inches(1), Inches(1.8), Inches(0.4),
             "키워드 매칭", font_size=10, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2.8), y_flow + Inches(1.3), Inches(1.8), Inches(0.5),
             "정확한 함수명,\n에러 문자열, 설정 키", font_size=9, color=BODY_WHITE, alignment=PP_ALIGN.CENTER)

# Below the vector path
add_glass_card(slide, Inches(0.5), Inches(4.2), Inches(5.5), Inches(2.8), alpha=10)
add_text_box(slide, Inches(0.7), Inches(4.35), Inches(5.1), Inches(0.4),
             "Vector (의미 검색)", font_size=14, color=ACCENT_VIOLET, bold=True)
add_bullet_text(slide, Inches(0.7), Inches(4.8), Inches(5.1), Inches(2),
                [
                    "• \"데이터베이스 성능 최적화\" → N+1 쿼리 수정 기록 찾음",
                    "• 단어가 달라도 의미가 같으면 매칭",
                    "• 8개 임베딩 프로바이더 (OpenAI, Gemini, Ollama 등)",
                    "• API 키만 있으면 자동 감지",
                ], font_size=11, color=BODY_WHITE)

add_glass_card(slide, Inches(6.5), Inches(4.2), Inches(5.5), Inches(2.8), alpha=10)
add_text_box(slide, Inches(6.7), Inches(4.35), Inches(5.1), Inches(0.4),
             "검색 품질 튜닝", font_size=14, color=ACCENT_CYAN, bold=True)
add_bullet_text(slide, Inches(6.7), Inches(4.8), Inches(5.1), Inches(2),
                [
                    "• 시간 감쇠 (Temporal Decay) — 최근 정보 우선",
                    "• MMR 다양성 — 중복 결과 제거",
                    "• CJK 트라이그램 — 한중일 자연어 처리",
                    "• 청크 분할 (~400토큰, 80토큰 오버랩)",
                ], font_size=11, color=BODY_WHITE)


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — QMD DEEP DIVE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_glow_circle(slide, Inches(9), Inches(2), Inches(4), Inches(4), GLOW_VIOLET, alpha=8)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.5),
             "BACKEND — QMD", font_size=11, color=ACCENT_VIOLET, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1.5), ACCENT_VIOLET)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.8),
             "QMD — 로컬 검색 사이드카 (by Tobi)",
             font_size=30, color=TITLE_WHITE, bold=True, font_name='Segoe UI Light')

# Left: What it adds
add_glass_card(slide, Inches(0.6), Inches(2.3), Inches(6), Inches(4.8), alpha=10)
add_text_box(slide, Inches(0.8), Inches(2.45), Inches(5.6), Inches(0.4),
             "Builtin 위에 추가되는 것", font_size=16, color=ACCENT_VIOLET, bold=True)
add_accent_line(slide, Inches(0.8), Inches(2.9), Inches(2), ACCENT_VIOLET)

items = [
    ("리랭킹 (Reranking)", "1차 검색 → 후보 추출 → 별도 모델로 재평가\n\"정확도가 높은 결과가 위로\""),
    ("쿼리 확장 (Query Expansion)", "\"인증\" → JWT, OAuth, 미들웨어까지 자동 확장\n연관 쿼리를 생성해서 더 넓게 검색"),
    ("외부 디렉토리 인덱싱", "~/notes, ~/projects 등 워크스페이스 밖도 검색\nBuiltin은 절대 못 하는 기능"),
    ("과거 대화 인덱싱", "세션 트랜스크립트를 검색 가능하게 변환\n\"저번 달에 뭐라고 했더라?\" 가능"),
    ("완전 로컬", "GGUF 모델 자동 다운로드 (~2GB)\nAPI 키 불필요, 인터넷 없이도 작동"),
]

y = Inches(3.1)
for title, desc in items:
    add_text_box(slide, Inches(0.8), y, Inches(5.6), Inches(0.3),
                 title, font_size=12, color=ACCENT_CYAN, bold=True)
    add_text_box(slide, Inches(0.8), y + Inches(0.3), Inches(5.6), Inches(0.5),
                 desc, font_size=10, color=BODY_WHITE)
    y += Inches(0.85)

# Right: When to use
add_glass_card(slide, Inches(7), Inches(2.3), Inches(5.7), Inches(4.8), alpha=10)
add_text_box(slide, Inches(7.2), Inches(2.45), Inches(5.3), Inches(0.4),
             "실전 의미", font_size=16, color=ACCENT_GREEN, bold=True)
add_accent_line(slide, Inches(7.2), Inches(2.9), Inches(2), ACCENT_GREEN)

add_bullet_text(slide, Inches(7.2), Inches(3.2), Inches(5.3), Inches(3.5),
                [
                    "✅ API 키 안 쓰고 완전 로컬로 검색하고 싶을 때",
                    "✅ 검색 정확도가 중요한 프로젝트",
                    "✅ 과거 대화 내용을 다시 찾아야 할 때",
                    "✅ 워크스페이스 밖 파일도 검색 범위에 포함",
                    "",
                    "⚠️ 첫 실행 시 GGUF 모델 다운로드로 느림",
                    "⚠️ 별도 바이너리 설치 필요 (npm/bun)",
                    "",
                    "🛡️ QMD가 죽으면 자동으로 Builtin 폴백",
                    "   → 에이전트는 계속 정상 작동",
                ], font_size=11, color=BODY_WHITE)


# ══════════════════════════════════════════════════════════════
# SLIDE 8 — HONCHO DEEP DIVE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_glow_circle(slide, Inches(-1), Inches(3), Inches(4), Inches(4), GLOW_CYAN, alpha=8)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.5),
             "BACKEND — HONCHO", font_size=11, color=ACCENT_GREEN, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1.5), ACCENT_GREEN)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.8),
             "Honcho — AI가 사용자를 이해하는 메모리",
             font_size=30, color=TITLE_WHITE, bold=True, font_name='Segoe UI Light')

# Core concept
add_glass_card(slide, Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.5), alpha=10)
add_text_box(slide, Inches(0.8), Inches(2.35), Inches(5.4), Inches(0.4),
             "근본적 차이", font_size=16, color=ACCENT_GREEN, bold=True)
add_text_box(slide, Inches(0.8), Inches(2.8), Inches(5.4), Inches(1.8),
             "다른 백엔드: 파일 → 청크 → 인덱스 (검색 중심)\n\n"
             "Honcho: 대화 전체를 관찰 → 사용자 모델링\n"
             "\"이 사람은 TypeScript를 선호하는구나\"를 자동 학습\n"
             "\"기억해\"라고 안 해도 알아서 기억",
             font_size=12, color=BODY_WHITE)

# Tools
add_glass_card(slide, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.5), alpha=10)
add_text_box(slide, Inches(7), Inches(2.35), Inches(5.4), Inches(0.4),
             "6개 도구", font_size=16, color=ACCENT_CYAN, bold=True)
add_bullet_text(slide, Inches(7), Inches(2.8), Inches(5.4), Inches(1.8),
                [
                    "빠른 조회 (LLM 없음):",
                    "  honcho_context / search_conclusions / search_messages / session",
                    "",
                    "LLM 기반 (더 깊은 이해):",
                    "  honcho_ask (quick=사실, thorough=종합분석)",
                ], font_size=11, color=BODY_WHITE)

# Bottom: Multi-agent + Comparison
add_glass_card(slide, Inches(0.6), Inches(5), Inches(5.8), Inches(2.3), alpha=10)
add_text_box(slide, Inches(0.8), Inches(5.15), Inches(5.4), Inches(0.4),
             "멀티 에이전트 인식", font_size=16, color=ACCENT_AMBER, bold=True)
add_bullet_text(slide, Inches(0.8), Inches(5.6), Inches(5.4), Inches(1.5),
                [
                    "• 부모 에이전트가 자식 에이전트를 자동 추적",
                    "• 채널/세션 경계 없이 동일한 사용자 컨텍스트",
                    "• Telegram, Discord, 웹 — 어디서든 같은 기억",
                    "• 로컬 셀프호스팅 또는 매니지드 API",
                ], font_size=11, color=BODY_WHITE)

add_glass_card(slide, Inches(6.8), Inches(5), Inches(5.8), Inches(2.3), alpha=10)
add_text_box(slide, Inches(7), Inches(5.15), Inches(5.4), Inches(0.4),
             "실전 의미", font_size=16, color=ACCENT_ROSE, bold=True)
add_bullet_text(slide, Inches(7), Inches(5.6), Inches(5.4), Inches(1.5),
                [
                    "✅ \"기억해\" 안 해도 자동으로 사용자를 이해",
                    "✅ 여러 채널에서 동일한 컨텍스트 유지",
                    "✅ 에이전트가 선호/스타일을 학습",
                    "⚠️ 플러그인 설치 + 별도 서버 필요",
                ], font_size=11, color=BODY_WHITE)


# ══════════════════════════════════════════════════════════════
# SLIDE 9 — LANCEDB DEEP DIVE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_glow_circle(slide, Inches(8), Inches(4), Inches(4), Inches(4), GLOW_VIOLET, alpha=6)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.5),
             "BACKEND — LANCEDB", font_size=11, color=ACCENT_AMBER, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1.5), ACCENT_AMBER)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.8),
             "LanceDB — 자동 리콜/캡처 로컬 벡터 DB",
             font_size=30, color=TITLE_WHITE, bold=True, font_name='Segoe UI Light')

# Core features
add_glass_card(slide, Inches(0.6), Inches(2.3), Inches(6), Inches(4.8), alpha=10)
add_text_box(slide, Inches(0.8), Inches(2.45), Inches(5.6), Inches(0.4),
             "핵심 기능", font_size=16, color=ACCENT_AMBER, bold=True)
add_accent_line(slide, Inches(0.8), Inches(2.9), Inches(2), ACCENT_AMBER)

items = [
    ("autoRecall", "매 턴 전에 사용자 메시지를 쿼리로 임베딩\n→ LanceDB에서 유사한 기억 찾기 → 컨텍스트에 주입"),
    ("autoCapture", "매 턴 후에 어시스턴트 응답이 기준 이하면\n→ 자동으로 LanceDB에 저장 (\"기억해\" 불필요)"),
    ("Ollama 네이티브", "/api/embed 직접 호출 — 완전 로컬, 무료\nmxbai-embed-large 등 로컬 모델 사용 가능"),
    ("S3 클라우드 저장", "dbPath를 s3://로 설정 → 여러 기기에서\n같은 메모리 공유 (멀티 디바이스 동기화)"),
]

y = Inches(3.1)
for title, desc in items:
    add_text_box(slide, Inches(0.8), y, Inches(1.8), Inches(0.3),
                 title, font_size=12, color=ACCENT_AMBER, bold=True)
    add_text_box(slide, Inches(2.8), y, Inches(3.6), Inches(0.6),
                 desc, font_size=10, color=BODY_WHITE)
    y += Inches(0.85)

# Right: Tools + Comparison
add_glass_card(slide, Inches(7), Inches(2.3), Inches(5.7), Inches(2.2), alpha=10)
add_text_box(slide, Inches(7.2), Inches(2.45), Inches(5.3), Inches(0.4),
             "에이전트 도구 3개", font_size=16, color=ACCENT_CYAN, bold=True)
add_bullet_text(slide, Inches(7.2), Inches(2.9), Inches(5.3), Inches(1.5),
                [
                    "memory_recall — 관련 기억 검색",
                    "memory_store — 사실/선호/결정/엔티티 저장",
                    "memory_forget — 매칭되는 기억 삭제",
                ], font_size=12, color=BODY_WHITE)

add_glass_card(slide, Inches(7), Inches(4.8), Inches(5.7), Inches(2.3), alpha=10)
add_text_box(slide, Inches(7.2), Inches(4.95), Inches(5.3), Inches(0.4),
             "Builtin vs LanceDB", font_size=16, color=ACCENT_VIOLET, bold=True)
add_bullet_text(slide, Inches(7.2), Inches(5.4), Inches(5.3), Inches(1.5),
                [
                    "Builtin: 수동 검색 (memory_search 직접 호출)",
                    "LanceDB: 자동 리콜 + 자동 캡처",
                    "Builtin: SQLite (단일 기기)",
                    "LanceDB: S3 저장 (멀티 기기 동기화)",
                ], font_size=11, color=BODY_WHITE)


# ══════════════════════════════════════════════════════════════
# SLIDE 10 — L3: ACTIVE MEMORY
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_glow_circle(slide, Inches(6), Inches(-1), Inches(5), Inches(5), GLOW_CYAN, alpha=8)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.5),
             "LAYER 3", font_size=11, color=ACCENT_GREEN, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1.5), ACCENT_GREEN)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.8),
             "Active Memory — 질문하기 전에 기억 찾기",
             font_size=30, color=TITLE_WHITE, bold=True, font_name='Segoe UI Light')

# Flow diagram
flow_steps = [
    ("사용자 질문", ACCENT_AMBER),
    ("차단 서브에이전트\n메모리 검색", ACCENT_VIOLET),
    ("관련 기억 발견?\nYES → 숨겨서 주입", ACCENT_GREEN),
    ("메인 에이전트\n자연스럽게 기억 활용", ACCENT_CYAN),
]

for i, (text, color) in enumerate(flow_steps):
    x = Inches(0.6) + i * Inches(3.15)
    y = Inches(2.3)
    add_glass_card(slide, x, y, Inches(2.8), Inches(1.4), alpha=10)
    add_icon_circle(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.5), color, str(i+1))
    add_text_box(slide, x + Inches(0.8), y + Inches(0.2), Inches(1.8), Inches(1),
                 text, font_size=11, color=TITLE_WHITE, bold=True, alignment=PP_ALIGN.LEFT)
    if i < len(flow_steps) - 1:
        add_text_box(slide, x + Inches(2.65), y + Inches(0.4), Inches(0.5), Inches(0.5),
                     "→", font_size=24, color=color, bold=True, alignment=PP_ALIGN.CENTER)

# Query modes
add_glass_card(slide, Inches(0.6), Inches(4.2), Inches(3.8), Inches(3), alpha=10)
add_text_box(slide, Inches(0.8), Inches(4.35), Inches(3.4), Inches(0.4),
             "쿼리 모드", font_size=16, color=ACCENT_CYAN, bold=True)
add_bullet_text(slide, Inches(0.8), Inches(4.8), Inches(3.4), Inches(2),
                [
                    "message — 마지막 메시지만 (빠름)",
                    "recent — 최근 대화 꼬리 포함 (균형)",
                    "full — 전체 대화 (정확, 느림)",
                ], font_size=11, color=BODY_WHITE)

# Prompt styles
add_glass_card(slide, Inches(4.7), Inches(4.2), Inches(3.8), Inches(3), alpha=10)
add_text_box(slide, Inches(4.9), Inches(4.35), Inches(3.4), Inches(0.4),
             "프롬프트 스타일", font_size=16, color=ACCENT_VIOLET, bold=True)
add_bullet_text(slide, Inches(4.9), Inches(4.8), Inches(3.4), Inches(2),
                [
                    "strict — 엄격하게",
                    "balanced — 균형 (기본값)",
                    "recall-heavy — 리콜 우선",
                    "preference-only — 선호만",
                ], font_size=11, color=BODY_WHITE)

# Key insight
add_glass_card(slide, Inches(8.8), Inches(4.2), Inches(3.8), Inches(3), alpha=10)
add_text_box(slide, Inches(9), Inches(4.35), Inches(3.4), Inches(0.4),
             "실전 의미", font_size=16, color=ACCENT_GREEN, bold=True)
add_bullet_text(slide, Inches(9), Inches(4.8), Inches(3.4), Inches(2),
                [
                    "✅ 에이전트가 기억을 먼저 찾음",
                    "✅ 사용자는 \"기억해\" 안 해도 됨",
                    "✅ 대화에 자연스럽게 녹아듦",
                    "⚠️ 전용 모델 권장 (속도)",
                ], font_size=11, color=BODY_WHITE)


# ══════════════════════════════════════════════════════════════
# SLIDE 11 — L4: DREAMING
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_glow_circle(slide, Inches(-1), Inches(5), Inches(4), Inches(4), GLOW_VIOLET, alpha=8)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.5),
             "LAYER 4", font_size=11, color=ACCENT_AMBER, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1.5), ACCENT_AMBER)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.8),
             "Dreaming — 백그라운드 메모리 통합",
             font_size=30, color=TITLE_WHITE, bold=True, font_name='Segoe UI Light')

# 3 phases
phases = [
    ("Light Sleep", "정리 / 스테이징", [
        "최근 단기 기억 정리",
        "후보 스테이징 (임시)",
        "강화 시그널 기록",
        "MEMORY.md에 안 씀",
    ], ACCENT_CYAN, "🌙"),
    ("Deep Sleep", "평가 / 승격", [
        "6개 시그널로 점수 매김",
        "3개 게이트 통과해야 승격",
        "MEMORY.md에 실제로 기록",
        "아무거나 다 저장하는 게 아님",
    ], ACCENT_VIOLET, "💤"),
    ("REM Sleep", "반영 / 메타", [
        "패턴과 주제 추출",
        "반영 시그널 기록",
        "Deep 랭킹에 피드백",
        "MEMORY.md에 안 씀",
    ], ACCENT_GREEN, "💭"),
]

for i, (title, subtitle, bullets, color, icon) in enumerate(phases):
    x = Inches(0.6) + i * Inches(4.15)
    y = Inches(2.3)
    add_glass_card(slide, x, y, Inches(3.9), Inches(3.5), alpha=10)
    add_icon_circle(slide, x + Inches(1.4), y + Inches(0.2), Inches(0.9), color, icon)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.3), Inches(3.5), Inches(0.4),
                 title, font_size=16, color=TITLE_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.7), Inches(3.5), Inches(0.3),
                 subtitle, font_size=10, color=color, alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, x + Inches(0.8), y + Inches(2.05), Inches(2), color)
    add_bullet_text(slide, x + Inches(0.3), y + Inches(2.2), Inches(3.3), Inches(1.2),
                    [f"• {b}" for b in bullets], font_size=10, color=BODY_WHITE)

# Ranking signals
add_glass_card(slide, Inches(0.6), Inches(6), Inches(7.5), Inches(1.3), alpha=10)
add_text_box(slide, Inches(0.8), Inches(6.1), Inches(7.1), Inches(0.4),
             "Deep 랭킹 6개 시그널", font_size=13, color=ACCENT_AMBER, bold=True)

signals = [
    ("빈도 24%", ACCENT_CYAN),
    ("관련성 30%", ACCENT_VIOLET),
    ("쿼리 다양성 15%", ACCENT_GREEN),
    ("최신성 15%", ACCENT_AMBER),
    ("통합 강도 10%", ACCENT_ROSE),
    ("개념 풍부도 6%", ACCENT_CYAN),
]

for i, (text, color) in enumerate(signals):
    x = Inches(0.8) + i * Inches(1.2)
    y = Inches(6.55)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.1), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.adjustments[0] = 0.2
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.color.rgb = BG_DEEP
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

# Right insight
add_glass_card(slide, Inches(8.5), Inches(6), Inches(4.3), Inches(1.3), alpha=10)
add_text_box(slide, Inches(8.7), Inches(6.1), Inches(3.9), Inches(0.3),
             "💡 기본 비활성화, 옵트인", font_size=12, color=ACCENT_AMBER, bold=True)
add_text_box(slide, Inches(8.7), Inches(6.45), Inches(3.9), Inches(0.7),
             "매일 새벽 3시 자동 실행\nDREAMS.md에서 사람이 리뷰 가능",
             font_size=10, color=BODY_WHITE)


# ══════════════════════════════════════════════════════════════
# SLIDE 12 — L5: CONTEXT ENGINE + COMPACTION
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_glow_circle(slide, Inches(7), Inches(2), Inches(5), Inches(5), GLOW_VIOLET, alpha=6)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.5),
             "LAYER 5", font_size=11, color=ACCENT_ROSE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1.5), ACCENT_ROSE)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.8),
             "Context Engine + Compaction",
             font_size=30, color=TITLE_WHITE, bold=True, font_name='Segoe UI Light')

# 4 lifecycle hooks
hooks = [
    ("ingest", "메시지 저장", "새 메시지가 들어오면\n엔진이 자체 저장소에 저장", ACCENT_CYAN),
    ("assemble", "컨텍스트 조립", "토큰 예산 내에서\n모델이 볼 메시지를 조립", ACCENT_VIOLET),
    ("compact", "압축", "대화가 길어지면\n오래된 내용을 요약", ACCENT_GREEN),
    ("afterTurn", "후처리", "실행 후 상태 업데이트\n백그라운드 작업 트리거", ACCENT_AMBER),
]

for i, (name, title, desc, color) in enumerate(hooks):
    x = Inches(0.6) + i * Inches(3.15)
    y = Inches(2.3)
    add_glass_card(slide, x, y, Inches(2.9), Inches(2.5), alpha=10)
    add_icon_circle(slide, x + Inches(0.9), y + Inches(0.2), Inches(0.8), color, name[:1].upper())
    add_text_box(slide, x + Inches(0.2), y + Inches(1.2), Inches(2.5), Inches(0.4),
                 f"{name}()", font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.6), Inches(2.5), Inches(0.3),
                 title, font_size=11, color=TITLE_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.95), Inches(2.5), Inches(0.5),
                 desc, font_size=9, color=BODY_WHITE, alignment=PP_ALIGN.CENTER)

# Compaction flow
add_glass_card(slide, Inches(0.6), Inches(5.2), Inches(5.8), Inches(2.1), alpha=10)
add_text_box(slide, Inches(0.8), Inches(5.35), Inches(5.4), Inches(0.4),
             "Compaction (자동 압축)", font_size=16, color=ACCENT_ROSE, bold=True)
add_bullet_text(slide, Inches(0.8), Inches(5.8), Inches(5.4), Inches(1.3),
                [
                    "1️⃣ 컨텍스트 한계 근접 → \"중요한 거 파일에 저장해\" 자동提醒",
                    "2️⃣ 오래된 대화를 요약으로 압축",
                    "3️⃣ 최근 메시지는 그대로 보존",
                    "4️⃣ 전체 기록은 디스크에 그대로 (검색 가능)",
                ], font_size=11, color=BODY_WHITE)

# Plugin engines
add_glass_card(slide, Inches(6.8), Inches(5.2), Inches(5.8), Inches(2.1), alpha=10)
add_text_box(slide, Inches(7), Inches(5.35), Inches(5.4), Inches(0.4),
             "플러그인 교체 가능", font_size=16, color=ACCENT_CYAN, bold=True)
add_bullet_text(slide, Inches(7), Inches(5.8), Inches(5.4), Inches(1.3),
                [
                    "기본: legacy 엔진 (OpenClaw 내장)",
                    "대체: lossless-claw 등 서드파티",
                    "컨텍스트 조립 방식을 완전히 커스텀 가능",
                    "DAG 요약, 벡터 검색 기반 등 자유로운 전략",
                ], font_size=11, color=BODY_WHITE)


# ══════════════════════════════════════════════════════════════
# SLIDE 13 — COMPARISON TABLE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_glow_circle(slide, Inches(5), Inches(3), Inches(5), Inches(5), GLOW_CYAN, alpha=5)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.5),
             "COMPARISON", font_size=11, color=ACCENT_CYAN, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1.5), ACCENT_CYAN)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.8),
             "4개 백엔드 비교",
             font_size=30, color=TITLE_WHITE, bold=True, font_name='Segoe UI Light')

# Table as glass cards
headers = ["", "Builtin", "QMD", "Honcho", "LanceDB"]
rows = [
    ["검색", "BM25+Vector", "+리랭크\n+쿼리확장", "의미 검색\n(관찰 기반)", "Vector"],
    ["자동 리콜", "❌", "❌", "✅ (프로파일)", "✅"],
    ["자동 캡처", "❌", "❌", "✅ (관찰)", "✅ (옵션)"],
    ["사용자 모델링", "❌", "❌", "✅ 자동", "❌"],
    ["과거 대화", "❌", "✅", "✅", "❌"],
    ["외부 디렉토리", "❌", "✅", "❌", "❌"],
    ["API 키", "필요", "불필요", "선택", "선택"],
    ["클라우드", "❌", "❌", "✅ (매니지드)", "✅ (S3)"],
]

col_width = Inches(2.4)
row_height = Inches(0.55)
start_x = Inches(0.6)
start_y = Inches(2.2)
header_height = Inches(0.5)

# Headers
for j, h in enumerate(headers):
    w = Inches(1.6) if j == 0 else col_width
    x = start_x + j * col_width if j > 0 else start_x
    if j > 0:
        x = start_x + Inches(1.6) + (j-1) * col_width
    add_glass_card(slide, x, start_y, w, header_height, alpha=15)
    add_text_box(slide, x + Inches(0.1), start_y + Inches(0.08), w - Inches(0.2), header_height - Inches(0.1),
                 h, font_size=12, color=ACCENT_CYAN if j > 0 else BODY_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Rows
for i, row in enumerate(rows):
    y = start_y + header_height + Inches(0.05) + i * (row_height + Inches(0.05))
    for j, cell in enumerate(row):
        w = Inches(1.6) if j == 0 else col_width
        x = start_x + j * col_width if j > 0 else start_x
        if j > 0:
            x = start_x + Inches(1.6) + (j-1) * col_width
        add_glass_card(slide, x, y, w, row_height, alpha=6)
        c = TITLE_WHITE if j == 0 else BODY_WHITE
        add_text_box(slide, x + Inches(0.05), y + Inches(0.05), w - Inches(0.1), row_height - Inches(0.1),
                     cell, font_size=9, color=c, bold=(j==0), alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 14 — PRACTICAL GUIDE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_glow_circle(slide, Inches(-1), Inches(-1), Inches(4), Inches(4), GLOW_VIOLET, alpha=8)
add_glow_circle(slide, Inches(10), Inches(5), Inches(4), Inches(4), GLOW_CYAN, alpha=8)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.5),
             "PRACTICAL GUIDE", font_size=11, color=ACCENT_GREEN, bold=True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1.5), ACCENT_GREEN)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.8),
             "어떤 조합을 쓸까?",
             font_size=30, color=TITLE_WHITE, bold=True, font_name='Segoe UI Light')

# 3 tiers
tiers = [
    ("🟢 기본 (0분 설정)", "Builtin만", [
        "API 키 하나로 하이브리드 검색 자동 활성",
        "MEMORY.md + 일일 노트로 충분",
        "대부분의 사용자가 여기서 만족",
    ], ACCENT_GREEN),
    ("🟡 중급 (10분 설정)", "Builtin + Active Memory", [
        "자동 리콜 활성화 → \"기억해\" 안 해도 됨",
        "빠른 모델 권장 (gemini-flash 등)",
        "개인 어시스턴트로 충분한 수준",
    ], ACCENT_AMBER),
    ("🔴 고급 (30분 설정)", "QMD/Honcho/LanceDB + Dreaming", [
        "QMD: 완전 로컬 + 과거 대화 검색",
        "Honcho: 자동 사용자 모델링",
        "LanceDB: 자동 리콜/캡처 + S3 동기화",
        "Dreaming: 백그라운드 통합",
    ], ACCENT_ROSE),
]

for i, (tier, subtitle, items, color) in enumerate(tiers):
    x = Inches(0.6) + i * Inches(4.15)
    y = Inches(2.3)
    add_glass_card(slide, x, y, Inches(3.9), Inches(4.5), alpha=10)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.3), Inches(3.5), Inches(0.5),
                 tier, font_size=16, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.8), Inches(3.5), Inches(0.3),
                 subtitle, font_size=11, color=BODY_WHITE)
    add_accent_line(slide, x + Inches(0.3), y + Inches(1.15), Inches(1.5), color)
    add_bullet_text(slide, x + Inches(0.3), y + Inches(1.4), Inches(3.3), Inches(2.8),
                    [f"• {b}" for b in items], font_size=11, color=BODY_WHITE, spacing=Pt(10))


# ══════════════════════════════════════════════════════════════
# SLIDE 15 — CLOSING
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
add_glow_circle(slide, Inches(2), Inches(1), Inches(5), Inches(5), GLOW_VIOLET, alpha=10)
add_glow_circle(slide, Inches(7), Inches(3), Inches(5), Inches(5), GLOW_CYAN, alpha=10)

card = add_glass_card(slide, Inches(2.5), Inches(1.5), Inches(8.333), Inches(4.5), alpha=12)

add_text_box(slide, Inches(3), Inches(2), Inches(7.333), Inches(1),
             "기억하는 에이전트, 기억하지 않는 에이전트",
             font_size=28, color=TITLE_WHITE, bold=True, alignment=PP_ALIGN.CENTER,
             font_name='Segoe UI Light')

add_accent_line(slide, Inches(5.5), Inches(3.1), Inches(2.333), ACCENT_CYAN)

add_text_box(slide, Inches(3), Inches(3.5), Inches(7.333), Inches(0.8),
             "차이는 사용자 경험 전부입니다",
             font_size=22, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(3), Inches(4.5), Inches(7.333), Inches(1.2),
             "OpenClaw는 5개 레이어로 이 차이를 메웁니다\n"
             "기본 설정만으로도 강력하고, 모두 켜면 인간처럼 기억합니다",
             font_size=14, color=BODY_WHITE, alignment=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────────
output_path = "/Users/conanssam-m4/.openclaw/workspace-blogbot/openclaw-memory-system.pptx"
prs.save(output_path)
print(f"✅ Saved: {output_path}")
