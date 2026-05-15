#!/usr/bin/env python3
"""OpenClaw Memory System PPTX — McKinsey / Corporate Consulting Style"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import lxml.etree as etree

# ── Constants ──────────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

# McKinsey / Swiss-Consulting palette
NAVY         = RGBColor(0x00, 0x2B, 0x5C)
DARK_NAVY    = RGBColor(0x00, 0x1A, 0x3A)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE    = RGBColor(0xF5, 0xF5, 0xF5)
LIGHT_GREY   = RGBColor(0xE8, 0xE8, 0xE8)
MED_GREY     = RGBColor(0x99, 0x99, 0x99)
DARK_GREY    = RGBColor(0x33, 0x33, 0x33)
BLACK        = RGBColor(0x1A, 0x1A, 0x1A)
SIGNAL_RED   = RGBColor(0xC8, 0x10, 0x2E)
GOLD         = RGBColor(0xB8, 0x86, 0x0B)
TEAL         = RGBColor(0x00, 0x7B, 0x7F)
SOFT_BLUE    = RGBColor(0x3A, 0x7C, 0xBD)
NAVY_LIGHT   = RGBColor(0x4A, 0x90, 0xD9)
AMBER        = RGBColor(0xD4, 0x8B, 0x0B)
GREEN_ACC    = RGBColor(0x1B, 0x7F, 0x4B)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
blank_layout = prs.slide_layouts[6]


def set_slide_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=DARK_GREY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri', italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = alignment
    return txBox

def add_rich_text(slide, left, top, width, height, lines, default_size=12,
                  default_color=DARK_GREY, default_font='Calibri'):
    """lines: list of (text, size, color, bold, alignment) tuples"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else default_color
        bold = line_data[3] if len(line_data) > 3 else False
        align = line_data[4] if len(line_data) > 4 else PP_ALIGN.LEFT
        italic = line_data[5] if len(line_data) > 5 else False
        
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = default_font
        p.font.italic = italic
        p.alignment = align
        p.space_after = Pt(4)
    return txBox

def add_rule_line(slide, left, top, width, color=SIGNAL_RED, thickness=Pt(2)):
    shape = add_rect(slide, left, top, width, thickness, color)
    return shape

def add_header_bar(slide, title_text):
    """McKinsey-style navy top bar with white title"""
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.9), NAVY)
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.6),
                 title_text, font_size=22, color=WHITE, bold=True,
                 font_name='Georgia')

def add_footer_bar(slide, page_num, total=15):
    """Subtle footer with page number"""
    add_rect(slide, Inches(0), H - Inches(0.3), W, Inches(0.3), LIGHT_GREY)
    add_text_box(slide, Inches(0.6), H - Inches(0.28), Inches(5), Inches(0.25),
                 "OpenClaw Memory System", font_size=8, color=MED_GREY, font_name='Calibri')
    add_text_box(slide, W - Inches(1.5), H - Inches(0.28), Inches(1), Inches(0.25),
                 f"{page_num}", font_size=8, color=MED_GREY, alignment=PP_ALIGN.RIGHT)

def add_divider(slide, left, top, height):
    """Vertical thin grey divider"""
    add_rect(slide, left, top, Pt(1), height, LIGHT_GREY)


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE (Navy full)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, NAVY)

# Thin gold accent line
add_rule_line(slide, Inches(1.2), Inches(2.2), Inches(2), GOLD, Pt(3))

add_text_box(slide, Inches(1.2), Inches(2.5), Inches(8), Inches(1.5),
             "OpenClaw\nMemory System",
             font_size=44, color=WHITE, bold=True, font_name='Georgia')

add_text_box(slide, Inches(1.2), Inches(4.3), Inches(8), Inches(0.6),
             "AI 에이전트가 기억하는 5개 레이어 아키텍처",
             font_size=18, color=RGBColor(0xAA, 0xCC, 0xE8), font_name='Calibri')

add_rule_line(slide, Inches(1.2), Inches(5.2), Inches(1.5), GOLD, Pt(2))

add_text_box(slide, Inches(1.2), Inches(5.5), Inches(8), Inches(0.5),
             "파일 기반 → 하이브리드 검색 → 자동 리콜 → 백그라운드 통합 → 컨텍스트 조립",
             font_size=11, color=RGBColor(0x88, 0xAA, 0xCC), font_name='Calibri')

add_text_box(slide, Inches(1.2), Inches(6.4), Inches(8), Inches(0.4),
             "2026.05  |  Concept & Architecture Review",
             font_size=10, color=MED_GREY, font_name='Calibri')


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_header_bar(slide, "문제 정의: AI 에이전트의 기억 상실")
add_footer_bar(slide, 2)

# Subtitle
add_text_box(slide, Inches(0.6), Inches(1.1), Inches(11), Inches(0.5),
             "현재 AI 코딩 에이전트의 근본적 한계 — 세션이 끝나면 모든 컨텍스트가 사라진다",
             font_size=14, color=DARK_GREY, font_name='Georgia', italic=True)

# 4 problem blocks in 2x2
problems = [
    ("01", "세션 리셋", "새 대화 시작 시 이전 컨텍스트 전부 소멸\nCLAUDE.md 200줄 한계로 정보 저장량 제한",
     "에이전트는 매 대화마다 깨끗한 상태에서 시작. 지난 세션에서 논의한 아키텍처, 결정사항, 버그 기록을 전혀 모름"),
    ("02", "수동 기억 의존", "\"기억해\"라고 직접 말해야 파일에 저장\n에이전트가 스스로 판단해 기억하지 못함",
     "사용자가 능동적으로 개입하지 않으면 아무것도 저장되지 않음. 인간의 장기 기억처럼 자동으로 작동하지 않음"),
    ("03", "검색 누락", "기억이 파일에 있어도 검색하지 못하면 무의미\n단순 키워드 매칭으로는 의미적 맥락을 못 찾음",
     "\"인증\"이라는 단어로 검색해도 JWT 설정, jose 선택 이유, 테스트 커버리지 같은 관련 맥락은 발견 불가"),
    ("04", "토큰 비효율", "전체 컨텍스트를 매번 복사하면 연간 1,950만 토큰\n비용 폭발 + 컨텍스트 윈도우 한계 초과",
     "LLM 윈도우 한계(200K) 때문에 전부 넣는 것도 불가능. 필요한 부분만 찾아서 넣는 메커니즘이 필수"),
]

for i, (num, title, desc, insight) in enumerate(problems):
    col = i % 2
    row = i // 2
    x = Inches(0.6) + col * Inches(6.2)
    y = Inches(1.8) + row * Inches(2.6)
    
    # Number circle
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.1), Inches(0.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SIGNAL_RED
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Title + desc
    add_text_box(slide, x + Inches(0.7), y, Inches(5.2), Inches(0.35),
                 title, font_size=15, color=NAVY, bold=True, font_name='Georgia')
    add_rule_line(slide, x + Inches(0.7), y + Inches(0.4), Inches(1.2), SIGNAL_RED, Pt(2))
    add_text_box(slide, x + Inches(0.7), y + Inches(0.55), Inches(5.2), Inches(0.6),
                 desc, font_size=10, color=DARK_GREY)
    
    # Insight box (light grey bg)
    add_rect(slide, x + Inches(0.7), y + Inches(1.3), Inches(5.2), Inches(1.0), OFF_WHITE, LIGHT_GREY, Pt(0.5))
    add_text_box(slide, x + Inches(0.85), y + Inches(1.4), Inches(4.9), Inches(0.8),
                 insight, font_size=9, color=MED_GREY, italic=True)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — 5 LAYER OVERVIEW
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_header_bar(slide, "아키텍처 개요: 5개 레이어 메모리 시스템")
add_footer_bar(slide, 3)

layers = [
    ("L1", "파일 기반 메모리", "MEMORY.md + 일일노트\n마크다운 파일 = 메모리", "모든 것의 기반. 복잡한 DB 없이 작동"),
    ("L2", "하이브리드 검색", "BM25 + Vector + RRF\n4종 백엔드 선택", "기억을 찾는 엔진. 의미 기반 검색 가능"),
    ("L3", "Active Memory", "자동 리콜 (질문 전)\n차단 서브에이전트", "\"기억해\" 없이도 자동으로 기억 활용"),
    ("L4", "Dreaming", "백그라운드 통합\nLight → Deep → REM", "인간의 수면 기억 정리를 모방"),
    ("L5", "Context Engine", "컨텍스트 조립\n+ Compaction", "모델이 볼 내용을 지능적으로 구성"),
]

# Vertical stack with connectors
for i, (num, title, desc, insight) in enumerate(layers):
    y = Inches(1.2) + i * Inches(1.15)
    
    # Number box
    shape = add_rect(slide, Inches(0.6), y, Inches(0.7), Inches(0.7), NAVY)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Layer content
    add_text_box(slide, Inches(1.5), y - Inches(0.05), Inches(2.5), Inches(0.35),
                 title, font_size=14, color=NAVY, bold=True, font_name='Georgia')
    add_rule_line(slide, Inches(1.5), y + Inches(0.3), Inches(1.5), SIGNAL_RED, Pt(1.5))
    add_text_box(slide, Inches(1.5), y + Inches(0.4), Inches(3), Inches(0.4),
                 desc, font_size=9, color=DARK_GREY)
    
    # Connector arrow
    if i < len(layers) - 1:
        add_text_box(slide, Inches(0.75), y + Inches(0.7), Inches(0.5), Inches(0.3),
                     "▼", font_size=10, color=NAVY, alignment=PP_ALIGN.CENTER)
    
    # Insight on the right
    add_rect(slide, Inches(5), y, Inches(7.8), Inches(0.7), OFF_WHITE, LIGHT_GREY, Pt(0.5))
    add_text_box(slide, Inches(5.2), y + Inches(0.1), Inches(7.4), Inches(0.5),
                 f"→ {insight}", font_size=10, color=DARK_GREY)


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — L1: FILE-BASED
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_header_bar(slide, "L1. 파일 기반 메모리 — 기본층")
add_footer_bar(slide, 4)

add_text_box(slide, Inches(0.6), Inches(1.1), Inches(11), Inches(0.4),
             "마크다운 파일 = 메모리. 열고, 읽고, 수정할 수 있는 투명한 저장소",
             font_size=13, color=DARK_GREY, font_name='Georgia', italic=True)

files_data = [
    ("MEMORY.md", "장기 기억", NAVY, [
        ("• 중요 사실, 결정, 선호를 큐레이션", 10, DARK_GREY, False),
        ("• 매 DM 세션 시작 시 자동 로딩", 10, DARK_GREY, False),
        ("• \"기억해\" → 자동으로 이 파일에 저장", 10, DARK_GREY, False),
        ("• 에이전트가 직접 편집하는 살아있는 문서", 10, DARK_GREY, False),
    ]),
    ("memory/YYYY-MM-DD.md", "일일 노트", TEAL, [
        ("• 당일 관찰, 작업 기록, 대화 요약", 10, DARK_GREY, False),
        ("• 오늘 + 어제 파일 자동 로딩", 10, DARK_GREY, False),
        ("• Dreaming의 원료 (단기 → 장기 승격)", 10, DARK_GREY, False),
        ("• 세션별로 자연 발생하는 기록", 10, DARK_GREY, False),
    ]),
    ("DREAMS.md", "꿈 일기 (옵션)", AMBER, [
        ("• Dreaming 시스템의 휴먼 리뷰용 출력", 10, DARK_GREY, False),
        ("• 승격 후보 + 반영 기록", 10, DARK_GREY, False),
        ("• 사람이 열어서 확인/수정 가능", 10, DARK_GREY, False),
        ("• 투명한 메모리 관리의 핵심", 10, DARK_GREY, False),
    ]),
]

for i, (fname, role, color, lines) in enumerate(files_data):
    x = Inches(0.6) + i * Inches(4.15)
    y = Inches(1.8)
    
    # Card background
    add_rect(slide, x, y, Inches(3.85), Inches(4.8), OFF_WHITE, LIGHT_GREY, Pt(0.5))
    # Color top bar
    add_rect(slide, x, y, Inches(3.85), Inches(0.06), color)
    
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.25), Inches(0.5),
                 fname, font_size=14, color=color, bold=True, font_name='Consolas')
    add_text_box(slide, x + Inches(0.3), y + Inches(0.7), Inches(3.25), Inches(0.3),
                 role, font_size=12, color=NAVY, bold=True, font_name='Georgia')
    add_rule_line(slide, x + Inches(0.3), y + Inches(1.05), Inches(1.2), color, Pt(1.5))
    
    add_rich_text(slide, x + Inches(0.3), y + Inches(1.3), Inches(3.25), Inches(3),
                  lines, default_font='Calibri')

# Bottom insight bar
add_rect(slide, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5), NAVY)
add_text_box(slide, Inches(0.9), Inches(6.75), Inches(11.5), Inches(0.4),
             "핵심:  마크다운 파일 자체가 메모리. DB 없이도 작동하며, 언제든 직접 열어 수정 가능",
             font_size=11, color=WHITE, bold=True)


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — L2: SEARCH OVERVIEW (4 backends)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_header_bar(slide, "L2. 검색 엔진 — 4종 백엔드 선택")
add_footer_bar(slide, 5)

add_text_box(slide, Inches(0.6), Inches(1.1), Inches(11), Inches(0.4),
             "기본값으로 충분하지만, 필요에 따라 더 강력한 백엔드로 교체 가능",
             font_size=13, color=DARK_GREY, font_name='Georgia', italic=True)

backends = [
    ("Builtin", "기본값", "SQLite + FTS5", NAVY, [
        "BM25 + Vector 하이브리드",
        "API 키 하나면 자동 활성",
        "CJK 트라이그램 지원",
        "별도 설치 불필요",
        "대부분의 사용자가 여기서 만족",
    ]),
    ("QMD", "로컬 고급", "리랭킹 + 쿼리 확장", TEAL, [
        "Builtin 위에 리랭킹 추가",
        "쿼리를 자동으로 확장",
        "외부 디렉토리 인덱싱",
        "과거 대화 검색 가능",
        "완전 로컬 (API 키 불필요)",
    ]),
    ("Honcho", "AI 네이티브", "사용자 모델링", SOFT_BLUE, [
        "자동 사용자 프로파일링",
        "멀티 에이전트 인식",
        "관찰(Observation) 기반 검색",
        "크로스 세션 / 크로스 채널",
        "로컬 또는 매니지드 API",
    ]),
    ("LanceDB", "벡터 DB", "자동 리콜/캡처", AMBER, [
        "LanceDB 컬럼 벡터 스토어",
        "autoRecall + autoCapture",
        "Ollama 임베딩 네이티브",
        "S3 클라우드 저장 지원",
        "멀티 디바이스 동기화",
    ]),
]

for i, (name, tag, subtitle, color, bullets) in enumerate(backends):
    x = Inches(0.5) + i * Inches(3.15)
    y = Inches(1.8)
    
    add_rect(slide, x, y, Inches(2.95), Inches(5), OFF_WHITE, LIGHT_GREY, Pt(0.5))
    add_rect(slide, x, y, Inches(2.95), Inches(0.06), color)
    
    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.55), Inches(0.4),
                 name, font_size=18, color=color, bold=True, font_name='Georgia')
    add_text_box(slide, x + Inches(0.2), y + Inches(0.6), Inches(2.55), Inches(0.25),
                 tag, font_size=9, color=MED_GREY, bold=True)
    add_rule_line(slide, x + Inches(0.2), y + Inches(0.95), Inches(1.2), color, Pt(1.5))
    add_text_box(slide, x + Inches(0.2), y + Inches(1.1), Inches(2.55), Inches(0.3),
                 subtitle, font_size=10, color=NAVY, bold=True)
    
    for j, bullet in enumerate(bullets):
        add_text_box(slide, x + Inches(0.2), y + Inches(1.5) + j * Inches(0.55), Inches(2.55), Inches(0.5),
                     f"• {bullet}", font_size=9, color=DARK_GREY)


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — HYBRID SEARCH DEEP DIVE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_header_bar(slide, "L2. 하이브리드 검색 원리")
add_footer_bar(slide, 6)

# Flow diagram
add_text_box(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.4),
             "두 개의 검색 경로를 병렬 실행 → RRF(Reciprocal Rank Fusion)로 병합",
             font_size=13, color=DARK_GREY, italic=True, font_name='Georgia')

# Horizontal flow
flow_steps = [
    ("Query\n(질문)", Inches(0.8), NAVY),
    ("Split", Inches(3.0), MED_GREY),
    ("BM25\n(키워드)", Inches(4.2), TEAL),
    ("Vector\n(의미)", Inches(4.2), SOFT_BLUE),
    ("RRF\nFusion", Inches(7.0), SIGNAL_RED),
    ("Top\nResults", Inches(9.2), GREEN_ACC),
]

y_flow = Inches(2.0)
# Query box
shape = add_rect(slide, Inches(0.6), y_flow, Inches(1.8), Inches(0.8), NAVY)
tf = shape.text_frame; p = tf.paragraphs[0]
p.text = "Query (질문)"; p.font.size = Pt(12); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER

# Arrow
add_text_box(slide, Inches(2.5), y_flow + Inches(0.15), Inches(0.5), Inches(0.5),
             "→", font_size=20, color=DARK_GREY, bold=True, alignment=PP_ALIGN.CENTER)

# Split point
add_text_box(slide, Inches(3.0), y_flow + Inches(0.15), Inches(0.8), Inches(0.5),
             "분기 ↓", font_size=11, color=MED_GREY, bold=True, alignment=PP_ALIGN.CENTER)

# BM25
shape = add_rect(slide, Inches(3.8), y_flow + Inches(0.9), Inches(2.2), Inches(0.7), TEAL)
tf = shape.text_frame; p = tf.paragraphs[0]
p.text = "BM25 (키워드)"; p.font.size = Pt(11); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER

# Vector
shape = add_rect(slide, Inches(3.8), y_flow + Inches(1.8), Inches(2.2), Inches(0.7), SOFT_BLUE)
tf = shape.text_frame; p = tf.paragraphs[0]
p.text = "Vector (의미)"; p.font.size = Pt(11); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER

# Merge arrow
add_text_box(slide, Inches(6.1), y_flow + Inches(1.1), Inches(0.5), Inches(0.5),
             "→", font_size=20, color=DARK_GREY, bold=True)
add_text_box(slide, Inches(6.1), y_flow + Inches(1.9), Inches(0.5), Inches(0.5),
             "→", font_size=20, color=DARK_GREY, bold=True)

# RRF
shape = add_rect(slide, Inches(6.7), y_flow + Inches(1.0), Inches(1.8), Inches(1.2), SIGNAL_RED)
tf = shape.text_frame; p = tf.paragraphs[0]
p.text = "RRF Fusion\n(병합)"; p.font.size = Pt(12); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER

# Arrow
add_text_box(slide, Inches(8.6), y_flow + Inches(1.2), Inches(0.5), Inches(0.5),
             "→", font_size=20, color=DARK_GREY, bold=True)

# Results
shape = add_rect(slide, Inches(9.2), y_flow + Inches(1.0), Inches(1.8), Inches(1.2), GREEN_ACC)
tf = shape.text_frame; p = tf.paragraphs[0]
p.text = "Top Results\n(결과)"; p.font.size = Pt(12); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER

# Descriptions
add_rect(slide, Inches(0.6), Inches(4.5), Inches(5.8), Inches(2.5), OFF_WHITE, LIGHT_GREY, Pt(0.5))
add_text_box(slide, Inches(0.8), Inches(4.6), Inches(5.4), Inches(0.3),
             "BM25 (키워드 매칭)", font_size=13, color=TEAL, bold=True, font_name='Georgia')
add_rule_line(slide, Inches(0.8), Inches(4.95), Inches(1.5), TEAL, Pt(1.5))
add_rich_text(slide, Inches(0.8), Inches(5.1), Inches(5.4), Inches(1.8),
    [
        ("• 정확한 함수명, 에러 문자열, 설정 키를 찾을 때 강력", 10, DARK_GREY, False),
        ("• SQLite FTS5 기반 — 추가 설치 불필요", 10, DARK_GREY, False),
        ("• CJK 트라이그램 토크나이저로 한글 자연어 처리", 10, DARK_GREY, False),
        ("\"auth_middleware\" → 정확히 그 파일을 찾음", 10, MED_GREY, False, PP_ALIGN.LEFT, True),
    ])

add_rect(slide, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.5), OFF_WHITE, LIGHT_GREY, Pt(0.5))
add_text_box(slide, Inches(7.0), Inches(4.6), Inches(5.4), Inches(0.3),
             "Vector (의미 검색)", font_size=13, color=SOFT_BLUE, bold=True, font_name='Georgia')
add_rule_line(slide, Inches(7.0), Inches(4.95), Inches(1.5), SOFT_BLUE, Pt(1.5))
add_rich_text(slide, Inches(7.0), Inches(5.1), Inches(5.4), Inches(1.8),
    [
        ("• 단어가 달라도 의미가 같으면 매칭", 10, DARK_GREY, False),
        ("• \"데이터베이스 성능\" → N+1 쿼리 수정 기록을 찾음", 10, DARK_GREY, False),
        ("• 8개 임베딩 프로바이더 자동 감지", 10, DARK_GREY, False),
        ("API 키만 있으면 자동 활성화 (OpenAI, Gemini, Ollama...)", 10, MED_GREY, False, PP_ALIGN.LEFT, True),
    ])


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — QMD DEEP DIVE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_header_bar(slide, "백엔드 심화: QMD — 로컬 검색 사이드카")
add_footer_bar(slide, 7)

# Two-column layout
# Left: Features
add_rect(slide, Inches(0.6), Inches(1.2), Inches(6), Inches(5.5), OFF_WHITE, LIGHT_GREY, Pt(0.5))
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.6), Inches(0.3),
             "Builtin 위에 추가되는 기능", font_size=15, color=TEAL, bold=True, font_name='Georgia')
add_rule_line(slide, Inches(0.8), Inches(1.65), Inches(2), TEAL, Pt(1.5))

qmd_features = [
    ("리랭킹 (Reranking)", "1차 검색 → 후보 추출 → 별도 리랭킹 모델로 재평가. 정확도가 높은 결과가 위로 올라옴"),
    ("쿼리 확장 (Query Expansion)", "\"인증\" 검색 → JWT, OAuth, 미들웨어까지 자동 확장. 연관 쿼리를 생성해서 더 넓게 검색"),
    ("외부 디렉토리 인덱싱", "~/notes, ~/projects 등 워크스페이스 밖의 마크다운 파일도 검색 범위에 포함"),
    ("과거 대화 인덱싱", "세션 트랜스크립트를 검색 가능하게 변환. \"저번 달에 뭐라고 했더라?\" 가능"),
    ("완전 로컬", "GGUF 모델 자동 다운로드 (~2GB). API 키 불필요, 인터넷 없이도 작동"),
]

y = Inches(1.9)
for title, desc in qmd_features:
    add_text_box(slide, Inches(0.8), y, Inches(5.6), Inches(0.3),
                 title, font_size=11, color=NAVY, bold=True)
    add_text_box(slide, Inches(0.8), y + Inches(0.3), Inches(5.6), Inches(0.5),
                 desc, font_size=9, color=DARK_GREY)
    y += Inches(0.9)

# Right: When to use
add_rect(slide, Inches(7), Inches(1.2), Inches(5.8), Inches(5.5), NAVY)
add_text_box(slide, Inches(7.2), Inches(1.4), Inches(5.4), Inches(0.3),
             "실전 의미 & 판단 기준", font_size=15, color=WHITE, bold=True, font_name='Georgia')
add_rule_line(slide, Inches(7.2), Inches(1.8), Inches(2), GOLD, Pt(1.5))

use_cases = [
    ("✅ API 키 없이 완전 로컬로 검색", WHITE),
    ("✅ 검색 정확도가 프로젝트 성패를 좌우", WHITE),
    ("✅ 과거 대화 기록을 재활용해야 할 때", WHITE),
    ("✅ 워크스페이스 밖 문서도 검색 필요", WHITE),
    ("", WHITE),
    ("⚠️ 첫 실행 시 GGUF 모델 다운로드로 느림", RGBColor(0xFF, 0xCC, 0x66)),
    ("⚠️ 별도 바이너리 설치 필요 (npm/bun)", RGBColor(0xFF, 0xCC, 0x66)),
    ("", WHITE),
    ("🛡️ 장애 대응: QMD가 죽으면 자동 Builtin 폴백", RGBColor(0x88, 0xDD, 0x88)),
    ("   에이전트는 계속 정상 작동", RGBColor(0x88, 0xDD, 0x88)),
]

y = Inches(2.1)
for text, color in use_cases:
    add_text_box(slide, Inches(7.2), y, Inches(5.4), Inches(0.3),
                 text, font_size=10, color=color)
    y += Inches(0.35)


# ══════════════════════════════════════════════════════════════
# SLIDE 8 — HONCHO DEEP DIVE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_header_bar(slide, "백엔드 심화: Honcho — AI 네이티브 사용자 기억")
add_footer_bar(slide, 8)

# Core difference
add_rect(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(1.2), NAVY)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.3),
             "근본적 차이: 파일 검색이 아니라 AI가 사용자를 모델링", font_size=14, color=WHITE, bold=True, font_name='Georgia')
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.5),
             "다른 백엔드: 텍스트 → 청크 → 인덱스 (검색 중심)  |  Honcho: 대화 전체 관찰 → 사용자 프로파일 자동 구축 → \"기억해\" 없이도 이해",
             font_size=11, color=RGBColor(0xAA, 0xCC, 0xE8))

# Tools (left)
add_rect(slide, Inches(0.6), Inches(2.7), Inches(6), Inches(2), OFF_WHITE, LIGHT_GREY, Pt(0.5))
add_text_box(slide, Inches(0.8), Inches(2.8), Inches(5.6), Inches(0.3),
             "6개 도구", font_size=13, color=SOFT_BLUE, bold=True, font_name='Georgia')
add_rule_line(slide, Inches(0.8), Inches(3.15), Inches(1.5), SOFT_BLUE, Pt(1.5))
add_rich_text(slide, Inches(0.8), Inches(3.3), Inches(5.6), Inches(1.3),
    [
        ("빠른 조회 (LLM 호출 없음):", 10, NAVY, True),
        ("  honcho_context / search_conclusions / search_messages / session", 9, DARK_GREY, False),
        ("", 6, DARK_GREY, False),
        ("LLM 기반 (더 깊은 이해):", 10, NAVY, True),
        ("  honcho_ask (quick=사실조회, thorough=종합분석)", 9, DARK_GREY, False),
    ])

# Multi-agent (right)
add_rect(slide, Inches(7), Inches(2.7), Inches(5.8), Inches(2), OFF_WHITE, LIGHT_GREY, Pt(0.5))
add_text_box(slide, Inches(7.2), Inches(2.8), Inches(5.4), Inches(0.3),
             "멀티 에이전트 인식", font_size=13, color=TEAL, bold=True, font_name='Georgia')
add_rule_line(slide, Inches(7.2), Inches(3.15), Inches(1.5), TEAL, Pt(1.5))
add_rich_text(slide, Inches(7.2), Inches(3.3), Inches(5.4), Inches(1.3),
    [
        ("• 부모 에이전트가 자식 에이전트를 자동 추적", 10, DARK_GREY, False),
        ("• 채널/세션 경계 없이 동일한 사용자 컨텍스트", 10, DARK_GREY, False),
        ("• Telegram, Discord, 웹 — 어디서든 같은 기억", 10, DARK_GREY, False),
        ("• 로컬 셀프호스팅 또는 매니지드 API (api.honcho.dev)", 10, DARK_GREY, False),
    ])

# Bottom: practical
add_rect(slide, Inches(0.6), Inches(5.0), Inches(12.1), Inches(2.2), NAVY)
add_text_box(slide, Inches(0.8), Inches(5.1), Inches(5.5), Inches(0.3),
             "실전 의미", font_size=13, color=GOLD, bold=True, font_name='Georgia')
add_rule_line(slide, Inches(0.8), Inches(5.45), Inches(1.5), GOLD, Pt(1.5))

add_rich_text(slide, Inches(0.8), Inches(5.6), Inches(5.5), Inches(1.5),
    [
        ("✅ \"기억해\" 없이도 사용자를 자동으로 이해", 10, WHITE, False),
        ("✅ 여러 채널에서 동일한 컨텍스트 유지", 10, WHITE, False),
        ("✅ 에이전트가 선호와 스타일을 학습", 10, WHITE, False),
        ("⚠️ 플러그인 설치 + 별도 서버 필요", 10, RGBColor(0xFF, 0xCC, 0x66), False),
    ])

add_text_box(slide, Inches(6.8), Inches(5.1), Inches(5.5), Inches(0.3),
             "작동 원리", font_size=13, color=GOLD, bold=True, font_name='Georgia')
add_rule_line(slide, Inches(6.8), Inches(5.45), Inches(1.5), GOLD, Pt(1.5))
add_rich_text(slide, Inches(6.8), Inches(5.6), Inches(5.5), Inches(1.5),
    [
        ("1. 매 턴 후: 대화 전체를 Honcho에 관찰로 저장", 10, WHITE, False),
        ("2. 사용자/에이전트 모델 자동 업데이트", 10, WHITE, False),
        ("3. 다음 세션: before_prompt_build에 관련 컨텍스트 주입", 10, WHITE, False),
        ("4. 에이전트가 자연스럽게 \"이 사람은 TS 선호\"를 앎", 10, WHITE, False),
    ])


# ══════════════════════════════════════════════════════════════
# SLIDE 9 — LANCEDB DEEP DIVE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_header_bar(slide, "백엔드 심화: LanceDB — 자동 리콜/캡처 로컬 벡터 DB")
add_footer_bar(slide, 9)

# Core features left
add_rect(slide, Inches(0.6), Inches(1.2), Inches(6), Inches(3), OFF_WHITE, LIGHT_GREY, Pt(0.5))
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.6), Inches(0.3),
             "핵심 기능", font_size=14, color=AMBER, bold=True, font_name='Georgia')
add_rule_line(slide, Inches(0.8), Inches(1.65), Inches(1.5), AMBER, Pt(1.5))

lancedb_features = [
    ("autoRecall (자동 리콜)", "매 턴 전에 사용자 메시지를 쿼리로 임베딩 → LanceDB에서 유사 기억 찾기 → 컨텍스트에 자동 주입"),
    ("autoCapture (자동 캡처)", "매 턴 후에 어시스턴트 응답이 기준 이하면 자동 저장. \"기억해\" 불필요"),
    ("Ollama 임베딩", "/api/embed 직접 호출 — mxbai-embed-large 등 로컬 모델. 완전 무료, API 키 불필요"),
    ("S3 클라우드 저장", "dbPath를 s3://로 설정 → 여러 기기에서 같은 메모리 공유 (멀티 디바이스 동기화)"),
]

y = Inches(1.85)
for title, desc in lancedb_features:
    add_text_box(slide, Inches(0.8), y, Inches(5.6), Inches(0.25),
                 title, font_size=10, color=NAVY, bold=True)
    add_text_box(slide, Inches(0.8), y + Inches(0.25), Inches(5.6), Inches(0.4),
                 desc, font_size=9, color=DARK_GREY)
    y += Inches(0.7)

# Tools + Comparison right
add_rect(slide, Inches(7), Inches(1.2), Inches(5.8), Inches(1.3), OFF_WHITE, LIGHT_GREY, Pt(0.5))
add_text_box(slide, Inches(7.2), Inches(1.3), Inches(5.4), Inches(0.3),
             "에이전트 도구 3개", font_size=13, color=NAVY, bold=True, font_name='Georgia')
add_rule_line(slide, Inches(7.2), Inches(1.65), Inches(1.5), NAVY, Pt(1.5))
add_rich_text(slide, Inches(7.2), Inches(1.8), Inches(5.4), Inches(0.6),
    [
        ("memory_recall — 관련 기억 검색  |  memory_store — 사실/선호 저장  |  memory_forget — 기억 삭제", 9, DARK_GREY, False),
    ])

# Builtin vs LanceDB
add_rect(slide, Inches(7), Inches(2.7), Inches(5.8), Inches(1.5), NAVY)
add_text_box(slide, Inches(7.2), Inches(2.8), Inches(5.4), Inches(0.3),
             "Builtin vs LanceDB", font_size=13, color=GOLD, bold=True, font_name='Georgia')
add_rich_text(slide, Inches(7.2), Inches(3.15), Inches(5.4), Inches(1),
    [
        ("Builtin: 수동 검색 (memory_search 직접 호출)", 9, WHITE, False),
        ("LanceDB: 자동 리콜 + 자동 캡처 — 설정만 하면 알아서 작동", 9, RGBColor(0x88, 0xDD, 0x88), False),
        ("Builtin: SQLite (단일 기기)  |  LanceDB: S3 저장 (멀티 기기)", 9, RGBColor(0xAA, 0xCC, 0xE8), False),
    ])

# Embedding options
add_rect(slide, Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.5), OFF_WHITE, LIGHT_GREY, Pt(0.5))
add_text_box(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(0.3),
             "임베딩 옵션 (다양한 환경 지원)", font_size=14, color=NAVY, bold=True, font_name='Georgia')
add_rule_line(slide, Inches(0.8), Inches(4.95), Inches(2), NAVY, Pt(1.5))

embed_options = [
    ("OpenAI", "text-embedding-3-small — API 키만 설정", TEAL),
    ("Ollama", "mxbai-embed-large — 완전 로컬, 무료", SOFT_BLUE),
    ("GitHub Copilot", "text-embedding-3-small — 구독으로 무료", GREEN_ACC),
    ("ZhiPu / 커스텀", "OpenAI 호환 엔드포인트 — dimensions 설정 필요", AMBER),
]

for i, (name, desc, color) in enumerate(embed_options):
    x = Inches(0.8) + i * Inches(3)
    add_rect(slide, x, Inches(5.2), Inches(2.8), Inches(1.5), WHITE, color, Pt(1.5))
    add_text_box(slide, x + Inches(0.15), Inches(5.3), Inches(2.5), Inches(0.3),
                 name, font_size=12, color=color, bold=True, font_name='Georgia')
    add_text_box(slide, x + Inches(0.15), Inches(5.65), Inches(2.5), Inches(0.8),
                 desc, font_size=9, color=DARK_GREY)


# ══════════════════════════════════════════════════════════════
# SLIDE 10 — ACTIVE MEMORY
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_header_bar(slide, "L3. Active Memory — 질문하기 전에 기억 찾기")
add_footer_bar(slide, 10)

add_text_box(slide, Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
             "대부분의 메모리 시스템은 수동 (에이전트가 검색 호출). Active Memory는 자동으로 먼저 찾는다",
             font_size=13, color=DARK_GREY, italic=True, font_name='Georgia')

# Flow: 4 steps horizontal
steps = [
    ("1", "사용자 질문 수신", NAVY),
    ("2", "차단 서브에이전트\n메모리 검색 실행", TEAL),
    ("3", "관련 기억 발견 시\n숨겨서 프롬프트에 주입", SOFT_BLUE),
    ("4", "메인 에이전트가\n자연스럽게 기억 활용", GREEN_ACC),
]

for i, (num, text, color) in enumerate(steps):
    x = Inches(0.6) + i * Inches(3.15)
    y = Inches(1.8)
    shape = add_rect(slide, x, y, Inches(2.8), Inches(1.2), color)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = f"Step {num}\n{text}"
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    if i < 3:
        add_text_box(slide, x + Inches(2.7), y + Inches(0.3), Inches(0.5), Inches(0.5),
                     "→", font_size=22, color=DARK_GREY, bold=True, alignment=PP_ALIGN.CENTER)

# Config options
add_rect(slide, Inches(0.6), Inches(3.3), Inches(4), Inches(3.5), OFF_WHITE, LIGHT_GREY, Pt(0.5))
add_text_box(slide, Inches(0.8), Inches(3.4), Inches(3.6), Inches(0.3),
             "쿼리 모드", font_size=13, color=NAVY, bold=True, font_name='Georgia')
add_rule_line(slide, Inches(0.8), Inches(3.75), Inches(1.2), NAVY, Pt(1.5))
add_rich_text(slide, Inches(0.8), Inches(3.9), Inches(3.6), Inches(2.5),
    [
        ("message", 11, NAVY, True),
        ("마지막 메시지만 전송. 가장 빠름", 9, DARK_GREY, False),
        ("", 6, DARK_GREY, False),
        ("recent  (권장)", 11, NAVY, True),
        ("최근 대화 꼬리 포함. 속도-정확도 균형", 9, DARK_GREY, False),
        ("", 6, DARK_GREY, False),
        ("full", 11, NAVY, True),
        ("전체 대화. 가장 정확하지만 느림", 9, DARK_GREY, False),
    ])

add_rect(slide, Inches(4.9), Inches(3.3), Inches(3.8), Inches(3.5), OFF_WHITE, LIGHT_GREY, Pt(0.5))
add_text_box(slide, Inches(5.1), Inches(3.4), Inches(3.4), Inches(0.3),
             "프롬프트 스타일", font_size=13, color=TEAL, bold=True, font_name='Georgia')
add_rule_line(slide, Inches(5.1), Inches(3.75), Inches(1.2), TEAL, Pt(1.5))
add_rich_text(slide, Inches(5.1), Inches(3.9), Inches(3.4), Inches(2.5),
    [
        ("balanced — 균형 (기본값)", 10, DARK_GREY, False),
        ("strict — 엄격하게 제한적", 10, DARK_GREY, False),
        ("recall-heavy — 리콜 최대화", 10, DARK_GREY, False),
        ("precision-heavy — 정밀도 우선", 10, DARK_GREY, False),
        ("preference-only — 선호만", 10, DARK_GREY, False),
        ("contextual — 맥락 강조", 10, DARK_GREY, False),
    ])

# Key insight
add_rect(slide, Inches(9), Inches(3.3), Inches(3.8), Inches(3.5), NAVY)
add_text_box(slide, Inches(9.2), Inches(3.5), Inches(3.4), Inches(0.3),
             "실전 의미", font_size=14, color=GOLD, bold=True, font_name='Georgia')
add_rule_line(slide, Inches(9.2), Inches(3.85), Inches(1.5), GOLD, Pt(1.5))
add_rich_text(slide, Inches(9.2), Inches(4.1), Inches(3.4), Inches(2.5),
    [
        ("✅ 에이전트가 먼저 기억을 찾음", 10, WHITE, False),
        ("✅ 사용자가 \"기억해\" 안 해도 됨", 10, WHITE, False),
        ("✅ 대화에 자연스럽게 녹아듦", 10, WHITE, False),
        ("✅ DM에만 기본 적용 (그룹 제외)", 10, WHITE, False),
        ("", 6, WHITE, False),
        ("권장: 빠른 모델 사용", 10, RGBColor(0xFF, 0xCC, 0x66), True),
        ("(gemini-flash / cerebras)", 9, RGBColor(0xFF, 0xCC, 0x66), False),
    ])


# ══════════════════════════════════════════════════════════════
# SLIDE 11 — DREAMING
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_header_bar(slide, "L4. Dreaming — 백그라운드 메모리 통합")
add_footer_bar(slide, 11)

add_text_box(slide, Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
             "인간의 수면 기억 정리를 모방. 단기 기억 → 장기 기억 자동 승격 (기본 비활성화)",
             font_size=13, color=DARK_GREY, italic=True, font_name='Georgia')

# 3 phases
phases = [
    ("Light Sleep", "정리 / 스테이징", TEAL, [
        "최근 단기 기억을 정리",
        "승격 후보를 스테이징 (임시)",
        "강화 시그널 기록",
        "MEMORY.md에 쓰지 않음 ✗",
    ]),
    ("Deep Sleep", "평가 / 승격", NAVY, [
        "6개 시그널로 가중 점수 계산",
        "3개 게이트 통과해야 승격",
        "MEMORY.md에 실제 기록 ✓",
        "아무거나 다 저장하지 않음",
    ]),
    ("REM Sleep", "반영 / 메타", AMBER, [
        "패턴과 주제를 추출",
        "반영 시그널 기록",
        "Deep 랭킹에 피드백 제공",
        "MEMORY.md에 쓰지 않음 ✗",
    ]),
]

for i, (title, subtitle, color, bullets) in enumerate(phases):
    x = Inches(0.6) + i * Inches(4.15)
    y = Inches(1.8)
    
    add_rect(slide, x, y, Inches(3.9), Inches(2.5), OFF_WHITE, LIGHT_GREY, Pt(0.5))
    add_rect(slide, x, y, Inches(3.9), Inches(0.06), color)
    
    add_text_box(slide, x + Inches(0.3), y + Inches(0.2), Inches(3.3), Inches(0.3),
                 title, font_size=15, color=color, bold=True, font_name='Georgia')
    add_text_box(slide, x + Inches(0.3), y + Inches(0.55), Inches(3.3), Inches(0.25),
                 subtitle, font_size=9, color=MED_GREY, bold=True)
    add_rule_line(slide, x + Inches(0.3), y + Inches(0.85), Inches(1.2), color, Pt(1.5))
    
    for j, bullet in enumerate(bullets):
        add_text_box(slide, x + Inches(0.3), y + Inches(1.0) + j * Inches(0.35), Inches(3.3), Inches(0.3),
                     f"• {bullet}", font_size=9, color=DARK_GREY)

# Ranking signals bar
add_rect(slide, Inches(0.6), Inches(4.6), Inches(12.1), Inches(0.5), NAVY)
add_text_box(slide, Inches(0.8), Inches(4.65), Inches(3), Inches(0.35),
             "Deep 랭킹 6개 시그널", font_size=11, color=WHITE, bold=True)

signals = [
    ("관련성 30%", Inches(4.0), TEAL),
    ("빈도 24%", Inches(5.8), SOFT_BLUE),
    ("쿼리 다양성 15%", Inches(7.3), GREEN_ACC),
    ("최신성 15%", Inches(9.1), AMBER),
    ("통합 강도 10%", Inches(10.5), SIGNAL_RED),
    ("개념 풍부도 6%", Inches(12.0), MED_GREY),
]

for text, x, color in signals:
    add_text_box(slide, x, Inches(4.68), Inches(1.5), Inches(0.3),
                 text, font_size=8, color=color, bold=True)

# Schedule + Key insight
add_rect(slide, Inches(0.6), Inches(5.3), Inches(6), Inches(1.8), OFF_WHITE, LIGHT_GREY, Pt(0.5))
add_text_box(slide, Inches(0.8), Inches(5.4), Inches(5.6), Inches(0.3),
             "스케줄 & 설정", font_size=13, color=NAVY, bold=True, font_name='Georgia')
add_rule_line(slide, Inches(0.8), Inches(5.75), Inches(1.2), NAVY, Pt(1.5))
add_rich_text(slide, Inches(0.8), Inches(5.9), Inches(5.6), Inches(1),
    [
        ("• 기본: 매일 새벽 3시 자동 실행 (cron)", 10, DARK_GREY, False),
        ("• DREAMS.md에서 사람이 리뷰 가능", 10, DARK_GREY, False),
        ("• 승격된 항목은 MEMORY.md에 추가", 10, DARK_GREY, False),
    ])

add_rect(slide, Inches(7), Inches(5.3), Inches(5.8), Inches(1.8), NAVY)
add_text_box(slide, Inches(7.2), Inches(5.4), Inches(5.4), Inches(0.3),
             "승격 조건 (3개 게이트 모두 통과해야)", font_size=12, color=GOLD, bold=True, font_name='Georgia')
add_rule_line(slide, Inches(7.2), Inches(5.75), Inches(2), GOLD, Pt(1.5))
add_rich_text(slide, Inches(7.2), Inches(5.9), Inches(5.4), Inches(1),
    [
        ("① 최소 점수 (minScore) — 충분히 중요한가?", 10, WHITE, False),
        ("② 최소 리콜 횟수 (minRecallCount) — 여러 번 참조되었나?", 10, WHITE, False),
        ("③ 최소 고유 쿼리 수 (minUniqueQueries) — 다양한 맥락에서 등장했나?", 10, WHITE, False),
    ])


# ══════════════════════════════════════════════════════════════
# SLIDE 12 — CONTEXT ENGINE + COMPACTION
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_header_bar(slide, "L5. Context Engine + Compaction")
add_footer_bar(slide, 12)

# 4 lifecycle hooks
hooks = [
    ("ingest()", "저장", "새 메시지가 들어오면\n엔진이 자체 저장소에 저장", NAVY),
    ("assemble()", "조립", "토큰 예산 내에서\n모델이 볼 메시지를 구성", TEAL),
    ("compact()", "압축", "대화가 길어지면\n오래된 내용을 요약", SOFT_BLUE),
    ("afterTurn()", "후처리", "실행 후 상태 업데이트\n백그라운드 작업 트리거", AMBER),
]

for i, (name, title, desc, color) in enumerate(hooks):
    x = Inches(0.6) + i * Inches(3.15)
    y = Inches(1.3)
    add_rect(slide, x, y, Inches(2.9), Inches(1.8), color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(2.5), Inches(0.3),
                 name, font_size=13, color=WHITE, bold=True, font_name='Consolas')
    add_text_box(slide, x + Inches(0.2), y + Inches(0.5), Inches(2.5), Inches(0.25),
                 title, font_size=10, color=RGBColor(0xCC, 0xDD, 0xEE), bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.85), Inches(2.5), Inches(0.8),
                 desc, font_size=9, color=WHITE)

# Compaction flow
add_rect(slide, Inches(0.6), Inches(3.4), Inches(6), Inches(3.5), OFF_WHITE, LIGHT_GREY, Pt(0.5))
add_text_box(slide, Inches(0.8), Inches(3.5), Inches(5.6), Inches(0.3),
             "Compaction (자동 압축) 프로세스", font_size=14, color=NAVY, bold=True, font_name='Georgia')
add_rule_line(slide, Inches(0.8), Inches(3.85), Inches(2), NAVY, Pt(1.5))

compaction_steps = [
    ("① 컨텍스트 한계 근접", "에이전트에게 \"중요한 거 파일에 저장해\"라고 자동 提醒 (Memory Flush)"),
    ("② 오래된 대화 요약", "LLM이 이전 대화를 요약 → 요약본으로 교체"),
    ("③ 최근 메시지 보존", "요약 이후의 최근 대화는 그대로 유지"),
    ("④ 전체 기록 보존", "디스크에는 전체 기록이 그대로 남음 → 검색 가능"),
]

y = Inches(4.1)
for title, desc in compaction_steps:
    add_text_box(slide, Inches(0.8), y, Inches(5.6), Inches(0.25),
                 title, font_size=10, color=NAVY, bold=True)
    add_text_box(slide, Inches(0.8), y + Inches(0.25), Inches(5.6), Inches(0.3),
                 desc, font_size=9, color=DARK_GREY)
    y += Inches(0.65)

# Plugin engines
add_rect(slide, Inches(7), Inches(3.4), Inches(5.8), Inches(3.5), NAVY)
add_text_box(slide, Inches(7.2), Inches(3.5), Inches(5.4), Inches(0.3),
             "플러그인으로 교체 가능", font_size=14, color=GOLD, bold=True, font_name='Georgia')
add_rule_line(slide, Inches(7.2), Inches(3.85), Inches(2), GOLD, Pt(1.5))

add_rich_text(slide, Inches(7.2), Inches(4.1), Inches(5.4), Inches(2.5),
    [
        ("기본 엔진: legacy", 11, WHITE, True),
        ("OpenClaw 내장 요약 기반 compaction", 10, RGBColor(0xAA, 0xCC, 0xE8), False),
        ("", 8, WHITE, False),
        ("서드파티 엔진 예시:", 11, WHITE, True),
        ("lossless-claw — 손실 없는 컨텍스트 관리", 10, RGBColor(0x88, 0xDD, 0x88), False),
        ("커스텀 엔진 — DAG 요약, 벡터 검색 기반 등", 10, RGBColor(0xAA, 0xCC, 0xE8), False),
        ("", 8, WHITE, False),
        ("컨텍스트 조립 방식을 완전히 커스텀 가능", 10, RGBColor(0xFF, 0xCC, 0x66), False),
        ("plugins.slots.contextEngine = \"my-engine\"", 9, RGBColor(0x88, 0xDD, 0x88), False),
    ])


# ══════════════════════════════════════════════════════════════
# SLIDE 13 — COMPARISON
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_header_bar(slide, "4개 백엔드 비교")
add_footer_bar(slide, 13)

# Table headers
cols = ["기능", "Builtin", "QMD", "Honcho", "LanceDB"]
col_w = [Inches(2.2), Inches(2.4), Inches(2.4), Inches(2.4), Inches(2.4)]
col_x = [Inches(0.6)]
for i in range(1, 5):
    col_x.append(col_x[i-1] + col_w[i-1] + Inches(0.05))

header_y = Inches(1.3)
for j, (h, x, w) in enumerate(zip(cols, col_x, col_w)):
    add_rect(slide, x, header_y, w, Inches(0.45), NAVY)
    add_text_box(slide, x + Inches(0.1), header_y + Inches(0.05), w - Inches(0.2), Inches(0.35),
                 h, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

rows = [
    ["검색", "BM25+Vector", "BM25+Vector\n+리랭크+쿼리확장", "의미 검색\n(관찰 기반)", "Vector"],
    ["자동 리콜", "✗", "✗", "✓ (프로파일)", "✓"],
    ["자동 캡처", "✗", "✗", "✓ (관찰)", "✓ (옵션)"],
    ["사용자 모델링", "✗", "✗", "✓ 자동", "✗"],
    ["과거 대화 검색", "✗", "✓", "✓", "✗"],
    ["외부 디렉토리", "✗", "✓", "✗", "✗"],
    ["API 키", "필요", "불필요", "선택", "선택"],
    ["클라우드 저장", "✗", "✗", "✓ (매니지드)", "✓ (S3)"],
    ["설치 난이도", "⭐", "⭐⭐", "⭐⭐", "⭐⭐"],
]

for i, row in enumerate(rows):
    y = header_y + Inches(0.55) + i * Inches(0.6)
    bg = OFF_WHITE if i % 2 == 0 else WHITE
    for j, (cell, x, w) in enumerate(zip(row, col_x, col_w)):
        add_rect(slide, x, y, w, Inches(0.5), bg, LIGHT_GREY, Pt(0.5))
        c = NAVY if j == 0 else DARK_GREY
        b = j == 0
        add_text_box(slide, x + Inches(0.1), y + Inches(0.05), w - Inches(0.2), Inches(0.4),
                     cell, font_size=9, color=c, bold=b, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 14 — PRACTICAL GUIDE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_header_bar(slide, "실전 가이드: 어떤 조합을 쓸까?")
add_footer_bar(slide, 14)

tiers = [
    ("기본", "0분 설정", "Builtin만", GREEN_ACC, [
        "API 키 하나로 하이브리드 검색 자동 활성",
        "MEMORY.md + 일일 노트로 충분",
        "대부분의 사용자가 여기서 만족",
    ]),
    ("중급", "10분 설정", "Builtin + Active Memory", AMBER, [
        "자동 리콜 활성화 → \"기억해\" 안 해도 됨",
        "빠른 모델 권장 (gemini-flash 등)",
        "개인 어시스턴트로 충분한 수준",
    ]),
    ("고급", "30분 설정", "QMD/Honcho/LanceDB + Dreaming", SIGNAL_RED, [
        "QMD: 완전 로컬 + 과거 대화 검색",
        "Honcho: 자동 사용자 모델링",
        "LanceDB: 자동 리콜/캡처 + S3 동기화",
        "Dreaming: 백그라운드 통합",
    ]),
]

for i, (tier, time, config, color, items) in enumerate(tiers):
    x = Inches(0.6) + i * Inches(4.15)
    y = Inches(1.3)
    
    add_rect(slide, x, y, Inches(3.9), Inches(5.5), OFF_WHITE, LIGHT_GREY, Pt(0.5))
    add_rect(slide, x, y, Inches(3.9), Inches(0.06), color)
    
    # Tier badge
    shape = add_rect(slide, x + Inches(0.3), y + Inches(0.3), Inches(1.2), Inches(0.4), color)
    tf = shape.text_frame; p = tf.paragraphs[0]
    p.text = tier; p.font.size = Pt(11); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    
    add_text_box(slide, x + Inches(1.7), y + Inches(0.3), Inches(1.8), Inches(0.4),
                 time, font_size=11, color=MED_GREY, bold=True)
    
    add_text_box(slide, x + Inches(0.3), y + Inches(0.9), Inches(3.3), Inches(0.3),
                 config, font_size=12, color=NAVY, bold=True, font_name='Georgia')
    add_rule_line(slide, x + Inches(0.3), y + Inches(1.25), Inches(1.5), color, Pt(1.5))
    
    for j, item in enumerate(items):
        add_text_box(slide, x + Inches(0.3), y + Inches(1.5) + j * Inches(0.55), Inches(3.3), Inches(0.5),
                     f"• {item}", font_size=10, color=DARK_GREY)

# Bottom recommendation
add_rect(slide, Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.4), NAVY)
add_text_box(slide, Inches(0.9), Inches(6.97), Inches(11.5), Inches(0.35),
             "권장:  기본에서 시작 → Active Memory 추가 → 필요시 고급 백엔드. 점진적 업그레이드가 핵심",
             font_size=11, color=WHITE, bold=True)


# ══════════════════════════════════════════════════════════════
# SLIDE 15 — CLOSING
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, NAVY)

add_rule_line(slide, Inches(1.2), Inches(2.5), Inches(2), GOLD, Pt(3))

add_text_box(slide, Inches(1.2), Inches(2.8), Inches(10), Inches(1.2),
             "기억하는 에이전트와\n기억하지 않는 에이전트",
             font_size=36, color=WHITE, bold=True, font_name='Georgia')

add_text_box(slide, Inches(1.2), Inches(4.3), Inches(10), Inches(0.5),
             "이 차이가 사용자 경험의 전부입니다",
             font_size=18, color=GOLD, font_name='Georgia')

add_rule_line(slide, Inches(1.2), Inches(5.0), Inches(1.5), GOLD, Pt(2))

add_text_box(slide, Inches(1.2), Inches(5.3), Inches(10), Inches(0.6),
             "OpenClaw는 5개 레이어로 이 차이를 메웁니다.\n기본 설정만으로도 강력하고, 모두 켜면 인간처럼 기억합니다.",
             font_size=13, color=RGBColor(0xAA, 0xCC, 0xE8))

add_text_box(slide, Inches(1.2), Inches(6.4), Inches(8), Inches(0.4),
             "2026.05  |  OpenClaw Memory System",
             font_size=10, color=MED_GREY, font_name='Calibri')


# ── Save ───────────────────────────────────────────────────
output_path = "/Users/conanssam-m4/Downloads/openclaw-memory-mckinsey.pptx"
prs.save(output_path)
print(f"✅ Saved: {output_path}")
